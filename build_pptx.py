"""
Field Manager Pro — Full 37-slide PowerPoint generator
Fixes: logo on section headers, corner icon on every slide, proper portrait phone mockups
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
import copy
from lxml import etree

# ── Paths ─────────────────────────────────────────────────────────────────────
LOGO_PATH = "/Users/shaungephart/Desktop/Claude Projects/field-manager-pro/ios/App/App/Assets.xcassets/AppIcon.appiconset/Icon-1024.png"
OUTPUT_PATH = "/Users/shaungephart/Desktop/Claude Projects/field-manager-pro/Field_Manager_Pro_User_Guide.pptx"

# ── Colors ─────────────────────────────────────────────────────────────────────
PURPLE      = RGBColor(0x6B, 0x21, 0xA8)
DARK_BG     = RGBColor(0x11, 0x18, 0x27)
MED_DARK    = RGBColor(0x1F, 0x29, 0x37)
CARD_BG     = RGBColor(0x37, 0x41, 0x51)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0x9C, 0xA3, 0xAF)
GREEN       = RGBColor(0x10, 0xB9, 0x81)
RED         = RGBColor(0xEF, 0x44, 0x44)
AMBER       = RGBColor(0xF5, 0x9E, 0x0B)
DARK_RED    = RGBColor(0x7F, 0x1D, 0x1D)
DARK_GREEN  = RGBColor(0x06, 0x4E, 0x3B)
DARK_AMBER  = RGBColor(0x78, 0x35, 0x00)
PHONE_BODY  = RGBColor(0x1A, 0x1A, 0x2E)
PHONE_SCREEN= RGBColor(0x11, 0x18, 0x27)

# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs

def blank_layout(prs):
    return prs.slide_layouts[6]  # completely blank

def fill_bg(slide, color):
    """Fill slide background with solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None,
             line_color=None, line_width_pt=0, corner_radius=None):
    """Add a rectangle (or rounded rect) shape."""
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    if corner_radius is not None:
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            left, top, width, height
        )
        # adjust[0] sets the corner radius as fraction (0=square, 0.5=max)
        shape.adjustments[0] = corner_radius
    else:
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            left, top, width, height
        )
    fill = shape.fill
    if fill_color:
        fill.solid()
        fill.fore_color.rgb = fill_color
    else:
        fill.background()
    line = shape.line
    if line_color:
        line.color.rgb = line_color
        if line_width_pt:
            line.width = Pt(line_width_pt)
    else:
        line.fill.background()
    return shape

def add_oval(slide, left, top, width, height, fill_color=None):
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        left, top, width, height
    )
    fill = shape.fill
    if fill_color:
        fill.solid()
        fill.fore_color.rgb = fill_color
    else:
        fill.background()
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=14,
                color=WHITE, bold=False, italic=False, align=PP_ALIGN.LEFT,
                word_wrap=True, transparency=0):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = word_wrap
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_para(tf, text, font_size=13, color=WHITE, bold=False, italic=False,
             align=PP_ALIGN.LEFT, space_before=0):
    from pptx.util import Pt as PPt
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = PPt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = PPt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p

def set_para_color(run, color):
    run.font.color.rgb = color

def add_logo(slide, x_in, y_in, size_in):
    """Add logo image."""
    slide.shapes.add_picture(
        LOGO_PATH,
        Inches(x_in), Inches(y_in),
        Inches(size_in), Inches(size_in)
    )

def add_corner_icon(slide):
    icon_w = Inches(1.575)
    icon_h = Inches(1.575)
    icon_x = Inches(11.505)
    icon_y = Inches(5.775)

    # Draw icon with no border
    slide.shapes.add_picture(LOGO_PATH, icon_x, icon_y, icon_w, icon_h)

    # Overlay a rounded rectangle with white border to match the icon's rounded corners
    overlay = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        icon_x, icon_y, icon_w, icon_h
    )
    overlay.adjustments[0] = 0.22  # matches iOS app icon corner radius (~22%)
    overlay.fill.background()       # fully transparent fill
    overlay.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    overlay.line.width = Pt(1.5)

def add_header_slide_text(slide, title, subtitle, footer):
    """Add consistently-positioned text around a truly-centered 6" logo on header slides.
    Logo occupies y=0.75" to y=6.75". Text fits in space above (0-0.75") and below (6.75-7.5").
    """
    # Title ABOVE logo: y=0.02", height=0.7"
    tb = slide.shapes.add_textbox(Inches(0.2), Inches(0.02), Inches(12.93), Inches(0.7))
    tb.word_wrap = False
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(52)
    run.font.bold = True
    run.font.color.rgb = WHITE

    # Subtitle BELOW logo: y=6.78", height=0.38"
    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(0.2), Inches(6.78), Inches(12.93), Inches(0.38))
        tb2.word_wrap = False
        tf2 = tb2.text_frame
        tf2.word_wrap = False
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(20)
        run2.font.color.rgb = WHITE
        # 80% opacity → set alpha via XML
        from lxml import etree
        from pptx.oxml.ns import qn
        solidFill = run2.font._element.find('.//' + qn('a:solidFill'))
        if solidFill is None:
            # force color set first then get element
            pass
        # Set transparency via srgbClr alpha
        r_xml = run2.font._element
        # Navigate to solidFill/srgbClr and add lumMod for opacity
        # Simpler: use theme color transparency not available easily; instead use 80% white approximation
        run2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xFF)

    # Footer: y=7.18", height=0.3"
    if footer:
        tb3 = slide.shapes.add_textbox(Inches(0.2), Inches(7.18), Inches(12.93), Inches(0.3))
        tb3.word_wrap = False
        tf3 = tb3.text_frame
        tf3.word_wrap = False
        p3 = tf3.paragraphs[0]
        p3.alignment = PP_ALIGN.CENTER
        run3 = p3.add_run()
        run3.text = footer
        run3.font.size = Pt(14)
        run3.font.color.rgb = RGBColor(0xBB, 0xBB, 0xEE)

# ── Phone Screenshot ──────────────────────────────────────────────────────────

SCREEN_DIR = "/Users/shaungephart/Desktop/Claude Projects/field-manager-pro/mockup_screens/"

def add_phone_screenshot(slide, image_path, left_inches=8.5, top_inches=1.2, height_inches=4.8):
    """Add a phone screenshot image to the slide, scaled to height_inches, maintaining aspect ratio."""
    if not os.path.exists(image_path):
        return
    # The screenshots are 280x580px (portrait). Aspect ratio = 280/580 = 0.4827
    aspect = 280 / 580
    h = Inches(height_inches)
    w = Inches(height_inches * aspect)
    left = Inches(left_inches)
    # Center vertically on slide (7.5" tall)
    top = Inches((7.5 - height_inches) / 2)
    slide.shapes.add_picture(image_path, left, top, w, h)

# ── Left-side content helpers ─────────────────────────────────────────────────

def left_heading(slide, text, y_in=0.55, x_in=0.5, w_in=6.5):
    add_textbox(slide, Inches(x_in), Inches(y_in), Inches(w_in), Inches(0.6),
                text, font_size=28, color=PURPLE, bold=True)

def left_body(slide, lines, start_y=1.35, x_in=0.5, w_in=6.5, size=13, color=WHITE, spacing=0.38):
    y = start_y
    for line in lines:
        add_textbox(slide, Inches(x_in), Inches(y), Inches(w_in), Inches(0.55),
                    line, font_size=size, color=color, word_wrap=True)
        y += spacing
    return y

def left_bullets(slide, items, start_y=1.35, x_in=0.5, w_in=6.5, size=13, color=WHITE, spacing=0.38):
    y = start_y
    for item in items:
        # purple dot
        add_oval(slide, Inches(x_in), Inches(y + 0.10), Inches(0.10), Inches(0.10),
                 fill_color=PURPLE)
        add_textbox(slide, Inches(x_in + 0.18), Inches(y), Inches(w_in - 0.2), Inches(0.5),
                    item, font_size=size, color=color, word_wrap=True)
        y += spacing
    return y

def section_header_slide(prs, title, subtitle="Your guide to using Field Manager Pro", footer=""):
    """Build a full-grape section header slide with large truly-centered logo."""
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, PURPLE)

    # Logo — truly centered: x=3.665", y=0.75"
    logo_size = 6.0
    logo_x = (13.33 - logo_size) / 2  # 3.665
    logo_y = 0.75
    add_logo(slide, logo_x, logo_y, logo_size)

    add_header_slide_text(slide, title=title, subtitle=subtitle, footer=footer)
    # No corner icon on section header slides
    return slide

def dark_card(slide, x_in, y_in, w_in, h_in, left_border_color=None):
    card = add_rect(slide, Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in),
                    fill_color=MED_DARK, corner_radius=0.05)
    if left_border_color:
        # Draw a thin colored left bar
        add_rect(slide, Inches(x_in), Inches(y_in), Inches(0.06), Inches(h_in),
                 fill_color=left_border_color)
    return card

# ══════════════════════════════════════════════════════════════════════════════
#  BUILD SLIDES
# ══════════════════════════════════════════════════════════════════════════════

def build_slide_01(prs):
    """Title Slide"""
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, PURPLE)

    # Logo — truly centered: x=3.665", y=0.75"
    logo_size = 6.0
    logo_x = (13.33 - logo_size) / 2  # 3.665
    logo_y = 0.75
    add_logo(slide, logo_x, logo_y, logo_size)

    add_header_slide_text(slide,
        title="Field Manager Pro",
        subtitle="User Guide & Role Responsibilities",
        footer="fieldmanagerpro.app")
    # No corner icon on title/header slides
    return slide

def build_slide_02(prs):
    """What is Field Manager Pro?"""
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, DARK_BG)

    left_heading(slide, "What is Field Manager Pro?")
    add_textbox(slide, Inches(0.5), Inches(1.25), Inches(6.5), Inches(0.7),
                "A mobile-first platform for managing field teams, schedules, tasks, and time tracking — all in one place.",
                font_size=13, color=LIGHT_GRAY, word_wrap=True)

    bullets = [
        "Clock In/Out Tracking",
        "Staff Scheduling",
        "Task Management",
        "Timecard Review",
        "Flags & Alerts",
        "Store Inspections",
        "Expense Tracking",
    ]
    left_bullets(slide, bullets, start_y=2.1, size=12, spacing=0.34)

    add_textbox(slide, Inches(0.5), Inches(6.6), Inches(6.5), Inches(0.5),
                "This guide is organized by role. Find your section and follow along.",
                font_size=12, color=LIGHT_GRAY, italic=True)

    # Phone screenshot — menu (overview)
    add_phone_screenshot(slide, SCREEN_DIR + "screen_menu.png", left_inches=8.7)

    add_corner_icon(slide)
    return slide

def build_slide_03(prs):
    """App Overview / Navigation"""
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, DARK_BG)

    left_heading(slide, "Accessing the App")
    steps = [
        "1. Open fieldmanagerpro.app on your phone or computer",
        "2. Sign in with your username and password",
        "3. Use the bottom navigation bar to switch between sections",
        "",
        "Forgot your password? Tap 'Forgot your password?' on the login screen.",
    ]
    left_body(slide, steps, start_y=1.3, size=13, spacing=0.42)

    # Phone screenshot — menu (app overview / navigation)
    add_phone_screenshot(slide, SCREEN_DIR + "screen_menu.png", left_inches=8.7)

    add_corner_icon(slide)
    return slide

def build_slide_04(prs):
    """General Rules"""
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, DARK_BG)

    add_textbox(slide, Inches(0.5), Inches(0.4), Inches(12.0), Inches(0.6),
                "General Rules — All Users", font_size=28, color=WHITE, bold=True)

    # Red warning box
    add_rect(slide, Inches(0.5), Inches(1.2), Inches(12.0), Inches(1.5),
             fill_color=DARK_RED, corner_radius=0.06)
    add_textbox(slide, Inches(0.7), Inches(1.28), Inches(11.6), Inches(0.38),
                "⚠  TIME FRAUD POLICY", font_size=16, color=WHITE, bold=True)
    add_textbox(slide, Inches(0.7), Inches(1.65), Inches(11.3), Inches(0.55),
                "Any attempt to share login credentials or clock in another employee will result in immediate termination.",
                font_size=13, color=WHITE, word_wrap=True)

    # Secondary box
    add_rect(slide, Inches(0.5), Inches(2.9), Inches(12.0), Inches(1.3),
             fill_color=MED_DARK, corner_radius=0.06)
    add_textbox(slide, Inches(0.7), Inches(3.0), Inches(11.3), Inches(1.1),
                "App Issues: If something isn't working, email your Owner with a description and screenshots of the error. Owner will escalate to the developer.",
                font_size=13, color=LIGHT_GRAY, word_wrap=True)

    # Third bullet
    add_rect(slide, Inches(0.5), Inches(4.4), Inches(12.0), Inches(0.7),
             fill_color=MED_DARK, corner_radius=0.06)
    add_oval(slide, Inches(0.72), Inches(4.58), Inches(0.10), Inches(0.10),
             fill_color=PURPLE)
    add_textbox(slide, Inches(0.9), Inches(4.5), Inches(11.0), Inches(0.4),
                "Never share your login credentials with anyone.", font_size=13, color=WHITE)

    add_corner_icon(slide)
    return slide

def build_slide_05(prs):
    """Employee Section Header"""
    return section_header_slide(prs,
        title="Employee",
        subtitle="Your guide to using Field Manager Pro",
        footer="")

def build_slide_06(prs):
    """Employee: Getting Started"""
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, DARK_BG)
    left_heading(slide, "Getting Started")

    steps = [
        "1. Go to fieldmanagerpro.app",
        "2. Enter your username and password",
        "3. Tap \"Sign In\"",
        "4. Use the bottom nav to explore your sections",
    ]
    y = left_body(slide, steps, start_y=1.3, size=13, spacing=0.40)
    add_textbox(slide, Inches(0.5), Inches(y + 0.15), Inches(6.5), Inches(0.8),
                "Your username and temporary password were sent to your email when your account was created.",
                font_size=12, color=LIGHT_GRAY, italic=True, word_wrap=True)

    # App store download section
    add_textbox(slide, Inches(0.5), Inches(5.0), Inches(6.5), Inches(0.35),
                "Available on iOS & Android",
                font_size=13, color=RGBColor(0x9C, 0xA3, 0xAF), italic=True)

    # App Store pill
    app_store_shape = add_rect(slide, Inches(0.5), Inches(5.3), Inches(1.6), Inches(0.35),
                               fill_color=RGBColor(0x1F, 0x29, 0x37), corner_radius=0.1)
    app_store_tf = app_store_shape.text_frame
    app_store_tf.word_wrap = False
    app_store_p = app_store_tf.paragraphs[0]
    app_store_p.alignment = PP_ALIGN.CENTER
    app_store_run = app_store_p.add_run()
    app_store_run.text = "\U0001F34E App Store"
    app_store_run.font.size = Pt(11)
    app_store_run.font.color.rgb = WHITE

    # Google Play pill
    gplay_shape = add_rect(slide, Inches(2.3), Inches(5.3), Inches(1.6), Inches(0.35),
                           fill_color=RGBColor(0x1F, 0x29, 0x37), corner_radius=0.1)
    gplay_tf = gplay_shape.text_frame
    gplay_tf.word_wrap = False
    gplay_p = gplay_tf.paragraphs[0]
    gplay_p.alignment = PP_ALIGN.CENTER
    gplay_run = gplay_p.add_run()
    gplay_run.text = "\u25B6 Google Play"
    gplay_run.font.size = Pt(11)
    gplay_run.font.color.rgb = WHITE

    add_textbox(slide, Inches(0.5), Inches(5.8), Inches(6.5), Inches(0.45),
                "Coming at launch \u2014 download for the best mobile experience",
                font_size=10, color=RGBColor(0x6B, 0x72, 0x80), italic=True, word_wrap=True)

    # Phone screenshot — login
    add_phone_screenshot(slide, SCREEN_DIR + "screen_login.png", left_inches=8.7)

    add_corner_icon(slide)
    return slide

def build_slide_07(prs):
    """Employee: Clocking In & Out"""
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, DARK_BG)
    left_heading(slide, "Clocking In & Out")

    lines = [
        "Tap the Clock icon in the navigation bar",
        "Tap Clock In when you arrive INSIDE the store",
        "and are ready to work",
        "Tap Clock Out BEFORE leaving the store",
    ]
    left_body(slide, lines, start_y=1.3, size=13, spacing=0.38)

    # Amber warning box
    add_rect(slide, Inches(0.5), Inches(3.5), Inches(6.5), Inches(1.4),
             fill_color=DARK_AMBER, corner_radius=0.06)
    add_textbox(slide, Inches(0.7), Inches(3.6), Inches(6.1), Inches(0.3),
                "IMPORTANT", font_size=13, color=AMBER, bold=True)
    add_textbox(slide, Inches(0.7), Inches(3.9), Inches(6.1), Inches(0.85),
                "Clock in inside the store after you are ready to work. Clock out before leaving. GPS location is recorded at both events.",
                font_size=12, color=WHITE, word_wrap=True)

    # Phone screenshot — clock in/out
    add_phone_screenshot(slide, SCREEN_DIR + "screen_clock.png", left_inches=8.7)

    add_corner_icon(slide)
    return slide

def build_slide_08(prs):
    """Employee: Viewing Your Schedule"""
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, DARK_BG)
    left_heading(slide, "Your Schedule")
    lines = [
        "Navigate to Schedule in the bottom navigation bar",
        "to view your upcoming shifts.",
        "",
        "Your DM posts the schedule. You will see your",
        "assigned days and hours.",
        "",
        "If you believe there's an error in your schedule,",
        "contact your DM.",
    ]
    left_body(slide, lines, start_y=1.3, size=13, spacing=0.34)

    # Phone screenshot — schedule
    add_phone_screenshot(slide, SCREEN_DIR + "screen_schedule.png", left_inches=8.7)

    add_corner_icon(slide)
    return slide

def build_slide_09(prs):
    """Employee: Checklist"""
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, DARK_BG)
    left_heading(slide, "Daily Checklist")
    lines = [
        "Navigate to Checklist to view and complete your daily assigned tasks.",
        "",
        "Tap each item to mark it complete.",
        "",
        "Complete all checklist items assigned to you each day.",
    ]
    left_body(slide, lines, start_y=1.3, size=13, spacing=0.38)

    # Phone screenshot — checklist
    add_phone_screenshot(slide, SCREEN_DIR + "screen_checklist.png", left_inches=8.7)

    add_corner_icon(slide)
    return slide

def build_slide_10(prs):
    """Employee: Time History"""
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, DARK_BG)
    left_heading(slide, "Time History")
    lines = [
        "Navigate to Time History to review your past",
        "clock-in and clock-out records.",
        "",
        "If you notice a discrepancy, notify your DM immediately.",
        "",
        "Your DM may make adjustments to your time entries",
        "and will notify you when they do.",
    ]
    left_body(slide, lines, start_y=1.3, size=13, spacing=0.35)

    # Phone screenshot — timecards
    add_phone_screenshot(slide, SCREEN_DIR + "screen_timecards.png", left_inches=8.7)

    add_corner_icon(slide)
    return slide

def build_slide_11(prs):
    """Employee: Expenses"""
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, DARK_BG)
    left_heading(slide, "Submitting Expenses")
    lines = [
        "Navigate to Expenses to submit a reimbursement request.",
        "",
        "1. Tap + New Expense",
        "2. Select Category",
        "3. Enter Amount",
        "4. Add description",
        "5. Attach receipt photo",
        "6. Submit — your DM will review",
    ]
    left_body(slide, lines, start_y=1.3, size=13, spacing=0.35)

    # Phone screenshot — expenses
    add_phone_screenshot(slide, SCREEN_DIR + "screen_expenses.png", left_inches=8.7)

    add_corner_icon(slide)
    return slide

def build_slide_12(prs):
    """DM Section Header"""
    return section_header_slide(prs,
        title="DM",
        subtitle="District Manager · Your guide to using Field Manager Pro",
        footer="")

def build_slide_13(prs):
    """DM: Overview & Responsibilities"""
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, DARK_BG)
    left_heading(slide, "DM Responsibilities", w_in=12.0)
    bullets = [
        "Manage your assigned team of employees",
        "Post the store schedule 2 WEEKS in advance at all times",
        "YOUR personal schedule must be posted 1 WEEK in advance — failure will prevent clock-in",
        "Review and validate timecards every Monday",
        "Check the app for new tasks and flags DAILY",
        "Clock in when you start your day; clock out only when your shift is fully complete",
        "You are responsible for your team's performance and schedule compliance",
    ]
    left_bullets(slide, bullets, start_y=1.25, x_in=0.5, w_in=12.2, size=13, spacing=0.62)
    add_corner_icon(slide)
    return slide

def build_slide_14(prs):
    """DM: Staff Scheduling"""
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, DARK_BG)
    left_heading(slide, "Creating the Staff Schedule")
    steps = [
        "1. Navigate to Staff Schedule",
        "2. Select the week using the arrows",
        "3. Tap + Add Shift",
        "4. Select employee, date, start time, end time",
        "5. Tap Publish to make schedule visible to your team",
    ]
    left_body(slide, steps, start_y=1.3, size=13, spacing=0.42)

    # Phone screenshot — staff schedule (staff scheduling)
    add_phone_screenshot(slide, SCREEN_DIR + "screen_staff_schedule.png", left_inches=8.7)

    add_corner_icon(slide)
    return slide

def build_slide_15(prs):
    """DM: Schedule Requirements"""
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, DARK_BG)
    left_heading(slide, "Schedule Requirements", w_in=7.5)

    # Req 1 — purple border
    dark_card(slide, 0.5, 1.3, 7.5, 1.30, left_border_color=PURPLE)
    add_textbox(slide, Inches(0.75), Inches(1.38), Inches(7.0), Inches(0.30),
                "Employee schedules must be posted 2 WEEKS in advance.", font_size=13, color=WHITE, bold=True)
    add_textbox(slide, Inches(0.75), Inches(1.68), Inches(7.0), Inches(0.50),
                "Example: Today is Monday, April 14 → Schedule must be complete through Sunday, April 27",
                font_size=11, color=LIGHT_GRAY, italic=True, word_wrap=True)

    # Req 2 — amber border
    dark_card(slide, 0.5, 2.80, 7.5, 1.30, left_border_color=AMBER)
    add_textbox(slide, Inches(0.75), Inches(2.88), Inches(7.0), Inches(0.30),
                "Your personal schedule must be posted 1 WEEK in advance.", font_size=13, color=WHITE, bold=True)
    add_textbox(slide, Inches(0.75), Inches(3.18), Inches(7.0), Inches(0.50),
                "Example: Today is Monday, April 14 → Your schedule must be set through Sunday, April 20",
                font_size=11, color=LIGHT_GRAY, italic=True, word_wrap=True)

    # Consequence — red border
    dark_card(slide, 0.5, 4.30, 7.5, 0.80, left_border_color=RED)
    add_textbox(slide, Inches(0.75), Inches(4.40), Inches(7.0), Inches(0.55),
                "⚠  If your personal schedule is not loaded, you CANNOT clock in.",
                font_size=13, color=RED, bold=True, word_wrap=True)

    # Phone screenshot — staff schedule (schedule requirements)
    add_phone_screenshot(slide, SCREEN_DIR + "screen_staff_schedule.png", left_inches=8.7)

    add_corner_icon(slide)
    return slide

def build_slide_16(prs):
    """DM: Timecard Review"""
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, DARK_BG)
    left_heading(slide, "Timecard Review")
    lines = [
        "Navigate to Timecards in the nav bar",
        "",
        "REQUIRED: Complete timecard review every Monday",
        "",
        "1. Review each employee's clock-in/out",
        "2. If an entry looks incorrect, tap it",
        "3. Edit the time",
        "4. Add a note explaining the change",
        "5. Save — employee is notified",
    ]
    left_body(slide, lines, start_y=1.3, size=13, spacing=0.35)

    # Phone screenshot — timecards
    add_phone_screenshot(slide, SCREEN_DIR + "screen_timecards.png", left_inches=8.7)

    add_corner_icon(slide)
    return slide

def build_slide_17(prs):
    """DM: Tasks"""
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, DARK_BG)
    left_heading(slide, "Managing Your Tasks")
    lines = [
        "Navigate to Tasks to view tasks assigned to you.",
        "",
        "Check for new tasks DAILY.",
        "",
        "To complete a task:",
        "  Tap the circle → Add a note (optional)",
        "  → Add a photo if required → Tap Mark Complete",
        "",
        "Tasks show a due date — complete them on time.",
    ]
    left_body(slide, lines, start_y=1.3, size=13, spacing=0.35)

    # Phone screenshot — tasks
    add_phone_screenshot(slide, SCREEN_DIR + "screen_tasks.png", left_inches=8.7)

    add_corner_icon(slide)
    return slide

def build_slide_18(prs):
    """DM: Flags"""
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, DARK_BG)
    left_heading(slide, "Reviewing Flags")
    lines = [
        "Flags are automatically created when unusual time activity is detected.",
        "",
        "Common flag types: Missed Clock-Out | Overtime Alert",
        "  | Late Clock-In | Missing Schedule",
        "",
        "Check flags DAILY. Resolve each flag with a note explaining what happened.",
        "",
        "1. Navigate to Flags",
        "2. Tap a flag to view details",
        "3. Enter resolution note",
        "4. Tap Resolve",
    ]
    left_body(slide, lines, start_y=1.3, size=13, spacing=0.35)

    # Phone screenshot — flags
    add_phone_screenshot(slide, SCREEN_DIR + "screen_flags.png", left_inches=8.7)

    add_corner_icon(slide)
    return slide

def build_slide_19(prs):
    """DM: Clocking In & Out"""
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, DARK_BG)
    left_heading(slide, "DM Clock In & Out")
    lines = [
        "As a DM, clock in when you START your workday.",
        "",
        "Clock out ONLY when your shift is completely finished for the day.",
        "",
        "GPS is recorded — your location is logged at clock-in and clock-out.",
        "",
        "Do not clock in for employees or allow others to clock in for you.",
    ]
    left_body(slide, lines, start_y=1.3, size=13, spacing=0.38)

    # Phone screenshot — clock in/out
    add_phone_screenshot(slide, SCREEN_DIR + "screen_clock.png", left_inches=8.7)

    add_corner_icon(slide)
    return slide

def build_slide_20(prs):
    """Ops Manager Section Header"""
    return section_header_slide(prs,
        title="Ops Manager",
        subtitle="Your guide to using Field Manager Pro",
        footer="")

def build_slide_21(prs):
    """Ops Manager: Responsibilities"""
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, DARK_BG)
    left_heading(slide, "Ops Manager Responsibilities", w_in=12.0)
    bullets = [
        "Tasks are assigned to you by the Owner",
        "Your access level is determined by what the Owner delegates to you",
        "Responsibilities vary — follow Owner's directives",
        "May manage scheduling, timecards, flags, and tasks depending on your assignment",
        "Apply DM-level standards to any area delegated to you",
        "Report issues and updates directly to the Owner",
    ]
    left_bullets(slide, bullets, start_y=1.3, x_in=0.5, w_in=12.2, size=13, spacing=0.62)
    add_corner_icon(slide)
    return slide

def build_slide_22(prs):
    """Ops Manager: Key Features"""
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, DARK_BG)
    left_heading(slide, "Your App Access")
    add_textbox(slide, Inches(0.5), Inches(1.25), Inches(6.5), Inches(0.4),
                "Depending on Owner delegation, you may have access to:",
                font_size=13, color=LIGHT_GRAY)
    features = [
        "Staff Schedule (view + manage)",
        "Timecards (view + edit)",
        "Tasks (receive + complete)",
        "Flags (view + resolve)",
        "Team (view)",
    ]
    y = 1.75
    for f in features:
        add_textbox(slide, Inches(0.5), Inches(y), Inches(0.3), Inches(0.30),
                    "✓", font_size=14, color=PURPLE, bold=True)
        add_textbox(slide, Inches(0.85), Inches(y), Inches(5.8), Inches(0.30),
                    f, font_size=13, color=WHITE)
        y += 0.42

    # Phone screenshot — menu (ops manager features)
    add_phone_screenshot(slide, SCREEN_DIR + "screen_menu.png", left_inches=8.7)

    add_corner_icon(slide)
    return slide

def build_slide_23(prs):
    """Sales Director Section Header"""
    return section_header_slide(prs,
        title="Sales Director",
        subtitle="Your guide to using Field Manager Pro",
        footer="")

def build_slide_24(prs):
    """Sales Director: Responsibilities"""
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, DARK_BG)
    left_heading(slide, "Sales Director Responsibilities", w_in=12.0)
    bullets = [
        "Assign and manage tasks for DMs",
        "Monitor schedule compliance — employee schedules must be 2 weeks in advance",
        "Review and validate timecards for your DMs",
        "Use the app as your primary management and inspection tool",
        "Ensure DMs are checking tasks and flags daily",
    ]
    left_bullets(slide, bullets, start_y=1.3, x_in=0.5, w_in=12.2, size=14, spacing=0.72)
    add_corner_icon(slide)
    return slide

def build_slide_25(prs):
    """Sales Director: Assigning Tasks"""
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, DARK_BG)
    left_heading(slide, "Assigning Tasks to DMs")
    steps = [
        "1. Navigate to Tasks",
        "2. Tap \"+ Assign Task\"",
        "3. Enter title and description",
        "4. Select assignee (DM or Ops Manager)",
        "5. Set a due date",
        "6. Tap \"Assign Task\" — assignee receives email",
        "",
        "Track completion: pending tasks show in the task list until marked complete.",
    ]
    left_body(slide, steps, start_y=1.3, size=13, spacing=0.37)

    # Phone screenshot — tasks (assigning tasks)
    add_phone_screenshot(slide, SCREEN_DIR + "screen_tasks.png", left_inches=8.7)

    add_corner_icon(slide)
    return slide

def build_slide_26(prs):
    """Sales Director: Schedule & Timecard Oversight"""
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, DARK_BG)
    left_heading(slide, "Schedule & Timecard Oversight")

    dark_card(slide, 0.5, 1.25, 7.5, 1.5, left_border_color=PURPLE)
    add_textbox(slide, Inches(0.75), Inches(1.33), Inches(7.0), Inches(0.28),
                "Schedule", font_size=14, color=PURPLE, bold=True)
    add_textbox(slide, Inches(0.75), Inches(1.63), Inches(7.0), Inches(0.80),
                "Navigate to Staff Schedule to view schedules across your stores. Verify DMs have schedules posted 2 weeks in advance.",
                font_size=13, color=WHITE, word_wrap=True)

    dark_card(slide, 0.5, 2.95, 7.5, 1.5, left_border_color=AMBER)
    add_textbox(slide, Inches(0.75), Inches(3.03), Inches(7.0), Inches(0.28),
                "Timecards", font_size=14, color=AMBER, bold=True)
    add_textbox(slide, Inches(0.75), Inches(3.33), Inches(7.0), Inches(0.80),
                "Navigate to Timecards to review DM time entries. Validate that hours are accurate and flag any discrepancies.",
                font_size=13, color=WHITE, word_wrap=True)

    # Phone screenshot — schedule (sales director oversight)
    add_phone_screenshot(slide, SCREEN_DIR + "screen_schedule.png", left_inches=8.7)

    add_corner_icon(slide)
    return slide

def build_slide_27(prs):
    """Owner Section Header"""
    return section_header_slide(prs,
        title="Owner",
        subtitle="Your guide to using Field Manager Pro",
        footer="")

def build_slide_28(prs):
    """Owner: Responsibilities"""
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, DARK_BG)
    left_heading(slide, "Owner Responsibilities", w_in=12.0)
    bullets = [
        "Primary operator of the organization in Field Manager Pro",
        "Run company operations and inspections through the app",
        "Assign tasks to Sales Director, Ops Manager, and DMs",
        "Oversee scheduling compliance, timecard accuracy, and flag resolution",
        "Manage your full team roster through the Team section",
        "First point of contact for app issues — forward to developer with screenshots",
    ]
    left_bullets(slide, bullets, start_y=1.3, x_in=0.5, w_in=12.2, size=13, spacing=0.65)
    add_corner_icon(slide)
    return slide

def build_slide_29(prs):
    """Owner: Team Management"""
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, DARK_BG)
    left_heading(slide, "Managing Your Team")
    steps = [
        "1. Navigate to Team",
        "2. Tap \"+ Add DM\" or \"+ Add Employee\"",
        "3. Fill in name, username, email, temporary password, and role",
        "4. Tap Create — user receives welcome email with login credentials",
        "",
        "Edit/deactivate: tap any user → Edit → adjust fields → Save Changes",
    ]
    left_body(slide, steps, start_y=1.3, size=13, spacing=0.40)

    # Phone screenshot — team management
    add_phone_screenshot(slide, SCREEN_DIR + "screen_team.png", left_inches=8.7)

    add_corner_icon(slide)
    return slide

def build_slide_30(prs):
    """Owner: Assigning Tasks"""
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, DARK_BG)
    left_heading(slide, "Assigning Tasks")
    lines = [
        "Navigate to Tasks → tap '+ Assign Task'",
        "",
        "Assign to: Sales Director, Ops Manager, or DM",
        "",
        "Set a due date for accountability — overdue tasks are highlighted in red",
        "",
        "Track completion: tasks show as Pending or Completed",
    ]
    left_body(slide, lines, start_y=1.3, size=13, spacing=0.42)

    # Phone screenshot — tasks (owner tasks)
    add_phone_screenshot(slide, SCREEN_DIR + "screen_tasks.png", left_inches=8.7)

    add_corner_icon(slide)
    return slide

def build_slide_31(prs):
    """Owner: DM Visit / Store Inspection"""
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, DARK_BG)
    left_heading(slide, "Store Inspections (DM Visit)")
    lines = [
        "Navigate to DM Visit to log store inspection visits.",
        "",
        "Record visit details, observations, and outcomes.",
        "",
        "Use this to track compliance and document findings.",
    ]
    left_body(slide, lines, start_y=1.3, size=13, spacing=0.42)

    # No screenshot available for DM Visit — slide left without phone mockup

    add_corner_icon(slide)
    return slide

def build_slide_32(prs):
    """Owner: Map View"""
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, DARK_BG)
    left_heading(slide, "Map View")
    lines = [
        "Navigate to Map to view store locations and GPS check-ins.",
        "",
        "See where employees clocked in and out.",
        "",
        "Use this to verify staff were on-site at time of clock-in.",
    ]
    left_body(slide, lines, start_y=1.3, size=13, spacing=0.45)

    # Phone screenshot — map view
    add_phone_screenshot(slide, SCREEN_DIR + "screen_map.png", left_inches=8.7)

    add_corner_icon(slide)
    return slide

def build_slide_33(prs):
    """RDM Section Header"""
    return section_header_slide(prs,
        title="RDM",
        subtitle="Regional Development Manager · Your guide to using Field Manager Pro",
        footer="")

def build_slide_34(prs):
    """RDM: Overview & Responsibilities"""
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, DARK_BG)
    left_heading(slide, "RDM Overview", w_in=12.0)
    add_textbox(slide, Inches(0.5), Inches(1.25), Inches(12.0), Inches(0.5),
                "The RDM role provides cross-organization visibility without administrative control.",
                font_size=13, color=LIGHT_GRAY, word_wrap=True)
    bullets = [
        "No clock in/out requirement",
        "View all staff schedules across all organizations",
        "Assign tasks to DMs across all organizations",
        "View-only access to schedules and task progress",
        "Report findings and task updates directly to the developer/platform owner",
    ]
    left_bullets(slide, bullets, start_y=1.85, x_in=0.5, w_in=12.2, size=13, spacing=0.62)
    add_corner_icon(slide)
    return slide

def build_slide_35(prs):
    """RDM: Assigning Tasks to DMs"""
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, DARK_BG)
    left_heading(slide, "Assigning Tasks to DMs")
    lines = [
        "Navigate to Tasks → tap '+ Assign Task'",
        "",
        "You can see DMs across all organizations in the assignee dropdown.",
        "",
        "Fill in title, description, and due date.",
        "",
        "The DM receives an email notification immediately.",
        "",
        "Monitor completion status in the Tasks view.",
    ]
    left_body(slide, lines, start_y=1.3, size=13, spacing=0.35)

    # Phone screenshot — tasks (RDM assigning tasks)
    add_phone_screenshot(slide, SCREEN_DIR + "screen_tasks.png", left_inches=8.7)

    add_corner_icon(slide)
    return slide

def build_slide_36(prs):
    """Support & Contact"""
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, DARK_BG)
    add_textbox(slide, Inches(0.5), Inches(0.4), Inches(12.0), Inches(0.6),
                "Support & App Issues", font_size=28, color=PURPLE, bold=True)

    # 3 steps
    steps_data = [
        ("Step 1", "Take a screenshot of the error message or issue"),
        ("Step 2", "Email your Owner with a description and the screenshot attached"),
        ("Step 3", "Owner will forward to the developer for resolution"),
    ]
    for i, (step, text) in enumerate(steps_data):
        sx2 = 0.5 + i * 4.1
        add_rect(slide, Inches(sx2), Inches(1.25), Inches(3.8), Inches(1.40),
                 fill_color=MED_DARK, corner_radius=0.06)
        add_textbox(slide, Inches(sx2 + 0.15), Inches(1.35), Inches(3.5), Inches(0.30),
                    step, font_size=12, color=PURPLE, bold=True)
        add_textbox(slide, Inches(sx2 + 0.15), Inches(1.68), Inches(3.5), Inches(0.80),
                    text, font_size=12, color=WHITE, word_wrap=True)

    # Red warning box
    add_rect(slide, Inches(0.5), Inches(2.95), Inches(12.33), Inches(1.30),
             fill_color=DARK_RED, corner_radius=0.06)
    add_textbox(slide, Inches(0.7), Inches(3.05), Inches(11.9), Inches(1.05),
                "Never share your login credentials. Do not attempt to clock in for another employee. Time fraud = immediate termination.",
                font_size=14, color=WHITE, bold=True, word_wrap=True)

    add_textbox(slide, Inches(0.5), Inches(4.45), Inches(12.0), Inches(0.35),
                "fieldmanagerpro.app", font_size=12, color=LIGHT_GRAY)

    add_corner_icon(slide)
    return slide

def build_slide_37(prs):
    """Welcome / Closing"""
    return section_header_slide(prs,
        title="Welcome to Field Manager Pro",
        subtitle="You now have everything you need to get started.",
        footer="fieldmanagerpro.app")

def build_support_slide(prs):
    """Support & Contact slide — appended at the end of each role section."""
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, DARK_BG)

    # Title
    add_textbox(slide, Inches(0.5), Inches(0.4), Inches(12.0), Inches(0.6),
                "Support & App Issues", font_size=28, color=PURPLE, bold=True)

    # Step boxes (dark cards)
    steps_data = [
        ("📸  Step 1: Take a screenshot of the error message or issue",),
        ("✉  Step 2: Email your Owner with a description and the screenshot attached",),
        ("📋  Step 3: Owner will forward to the developer for resolution",),
    ]
    for i, (text,) in enumerate(steps_data):
        sx2 = 0.5 + i * 4.1
        add_rect(slide, Inches(sx2), Inches(1.25), Inches(3.8), Inches(1.40),
                 fill_color=MED_DARK, corner_radius=0.06)
        add_textbox(slide, Inches(sx2 + 0.15), Inches(1.35), Inches(3.5), Inches(1.05),
                    text, font_size=12, color=WHITE, word_wrap=True)

    # Red warning box
    add_rect(slide, Inches(0.5), Inches(2.95), Inches(12.33), Inches(1.30),
             fill_color=DARK_RED, corner_radius=0.06)
    add_textbox(slide, Inches(0.7), Inches(3.05), Inches(11.9), Inches(1.05),
                "⚠ Never share your login credentials. Do not clock in for another employee. Time fraud = immediate termination.",
                font_size=14, color=WHITE, bold=True, word_wrap=True)

    # Small gray text at bottom
    add_textbox(slide, Inches(0.5), Inches(4.45), Inches(12.0), Inches(0.35),
                "fieldmanagerpro.app", font_size=12, color=LIGHT_GRAY)

    # Corner icon with white 1.5pt border
    add_corner_icon(slide)
    return slide


# ══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()

    # Slides 1–4: Intro / General
    for builder in [build_slide_01, build_slide_02, build_slide_03, build_slide_04]:
        print(f"  Building {builder.__name__} ...", flush=True)
        builder(prs)

    # Slides 5–11: Employee section
    for builder in [build_slide_05, build_slide_06, build_slide_07, build_slide_08,
                    build_slide_09, build_slide_10, build_slide_11]:
        print(f"  Building {builder.__name__} ...", flush=True)
        builder(prs)
    # Support slide after Employee section
    print("  Building support slide (after Employee) ...", flush=True)
    build_support_slide(prs)

    # Slides 12–19: DM section
    for builder in [build_slide_12, build_slide_13, build_slide_14, build_slide_15,
                    build_slide_16, build_slide_17, build_slide_18, build_slide_19]:
        print(f"  Building {builder.__name__} ...", flush=True)
        builder(prs)
    # Support slide after DM section
    print("  Building support slide (after DM) ...", flush=True)
    build_support_slide(prs)

    # Slides 20–22: Ops Manager section
    for builder in [build_slide_20, build_slide_21, build_slide_22]:
        print(f"  Building {builder.__name__} ...", flush=True)
        builder(prs)
    # Support slide after Ops Manager section
    print("  Building support slide (after Ops Manager) ...", flush=True)
    build_support_slide(prs)

    # Slides 23–26: Sales Director section
    for builder in [build_slide_23, build_slide_24, build_slide_25, build_slide_26]:
        print(f"  Building {builder.__name__} ...", flush=True)
        builder(prs)
    # Support slide after Sales Director section
    print("  Building support slide (after Sales Director) ...", flush=True)
    build_support_slide(prs)

    # Slides 27–32: Owner section
    for builder in [build_slide_27, build_slide_28, build_slide_29, build_slide_30,
                    build_slide_31, build_slide_32]:
        print(f"  Building {builder.__name__} ...", flush=True)
        builder(prs)
    # Support slide after Owner section
    print("  Building support slide (after Owner) ...", flush=True)
    build_support_slide(prs)

    # Slides 33–35: RDM section
    for builder in [build_slide_33, build_slide_34, build_slide_35]:
        print(f"  Building {builder.__name__} ...", flush=True)
        builder(prs)
    # Support slide after RDM section
    print("  Building support slide (after RDM) ...", flush=True)
    build_support_slide(prs)

    # Slides 36–37: Existing standalone support + Closing
    for builder in [build_slide_36, build_slide_37]:
        print(f"  Building {builder.__name__} ...", flush=True)
        builder(prs)

    prs.save(OUTPUT_PATH)
    size = os.path.getsize(OUTPUT_PATH)
    print(f"\nSaved: {OUTPUT_PATH}")
    print(f"File size: {size:,} bytes ({size/1024:.1f} KB)")
    print(f"Slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
