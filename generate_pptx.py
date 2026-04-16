from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.dml.color import RGBColor
import copy
from lxml import etree

# ── Color palette ──────────────────────────────────────────────────────────────
PURPLE       = RGBColor(0x7c, 0x3a, 0xed)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG      = RGBColor(0x11, 0x18, 0x27)
DARK2        = RGBColor(0x1f, 0x29, 0x37)
MED_GRAY     = RGBColor(0x37, 0x41, 0x51)
LIGHT_GRAY   = RGBColor(0x9c, 0xa3, 0xaf)
GREEN        = RGBColor(0x10, 0xb9, 0x81)
RED          = RGBColor(0xef, 0x44, 0x44)
ORANGE       = RGBColor(0xf5, 0x9e, 0x0b)
PHONE_SCREEN = RGBColor(0x11, 0x18, 0x27)

# ── Slide size: widescreen 13.33 x 7.5 ──────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Blank layout ──────────────────────────────────────────────────────────────
blank_layout = prs.slide_layouts[6]

slide_counter = [0]  # mutable for closure

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# HELPER UTILITIES
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def new_slide():
    slide_counter[0] += 1
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill=None, line=None, radius=False):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    if radius:
        # add rounded corners via XML
        sp = shape.element
        spPr = sp.find(qn('p:spPr'))
        prstGeom = spPr.find(qn('a:prstGeom'))
        if prstGeom is not None:
            prstGeom.set('prst', 'roundRect')
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is None:
                avLst = etree.SubElement(prstGeom, qn('a:avLst'))
            gd = etree.SubElement(avLst, qn('a:gd'))
            gd.set('name', 'adj')
            gd.set('fmla', 'val 30000')
    return shape


def add_textbox(slide, text, l, t, w, h,
                font_size=14, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def add_label_in_rect(slide, shape, text, font_size=11, bold=False,
                       color=WHITE, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_slide_number(slide, num):
    tb = slide.shapes.add_textbox(Inches(12.5), Inches(7.1), Inches(0.7), Inches(0.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(num)
    run.font.size = Pt(10)
    run.font.color.rgb = LIGHT_GRAY


def heading(slide, text, t=0.3, font_size=32, color=WHITE):
    add_textbox(slide, text, 0.4, t, 12.5, 0.8,
                font_size=font_size, bold=True, color=color, align=PP_ALIGN.LEFT)


def subheading(slide, text, t=1.0, color=LIGHT_GRAY, font_size=16):
    add_textbox(slide, text, 0.4, t, 12.5, 0.5,
                font_size=font_size, color=color, align=PP_ALIGN.LEFT)


def divider(slide, t=1.2, color=PURPLE, l=0.4, w=12.5):
    bar = add_rect(slide, l, t, w, 0.04, fill=color)
    return bar


def bullet_list(slide, items, l=0.5, t=1.5, w=5.8, font_size=14,
                color=WHITE, spacing=0.42, bullet_char="•"):
    for i, item in enumerate(items):
        y = t + i * spacing
        add_textbox(slide, f"{bullet_char}  {item}", l, y, w, 0.4,
                    font_size=font_size, color=color)


def section_header_slide(title, subtitle=""):
    slide = new_slide()
    fill_bg(slide, DARK2)
    # big purple block
    add_rect(slide, 0, 0, 13.33, 7.5, fill=PURPLE)
    # accent bar
    add_rect(slide, 0, 5.8, 13.33, 1.7, fill=DARK2)
    add_textbox(slide, title, 0.5, 2.5, 12.3, 1.4,
                font_size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        add_textbox(slide, subtitle, 0.5, 4.1, 12.3, 0.6,
                    font_size=20, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "fieldmanagerpro.app", 0.5, 6.1, 12.3, 0.5,
                font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    return slide


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# PHONE MOCKUP HELPER
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def phone_frame(slide, l=7.5, t=0.7, w=4.8, h=5.9):
    """Draw a phone-shaped container; returns (l,t,w,h) of inner screen."""
    # outer shell (slightly lighter)
    outer = add_rect(slide, l - 0.08, t - 0.08, w + 0.16, h + 0.16,
                     fill=MED_GRAY, radius=True)
    # screen
    screen = add_rect(slide, l, t, w, h, fill=PHONE_SCREEN, radius=True)
    # notch/status bar
    notch = add_rect(slide, l + w/2 - 0.4, t + 0.08, 0.8, 0.15,
                     fill=DARK2, radius=True)
    return (l, t, w, h)


def nav_bar(slide, l, t, w, h, labels):
    """Draw a bottom nav bar inside a phone mockup."""
    bar = add_rect(slide, l, t, w, h, fill=MED_GRAY)
    n = len(labels)
    icon_w = w / n
    for i, lbl in enumerate(labels):
        ix = l + i * icon_w + icon_w/2 - 0.25
        # circle icon
        circ = slide.shapes.add_shape(
            9,  # oval
            Inches(ix), Inches(t + 0.05), Inches(0.5), Inches(0.5)
        )
        circ.fill.solid()
        circ.fill.fore_color.rgb = PURPLE
        circ.line.fill.background()
        # label
        add_textbox(slide, lbl, ix - 0.1, t + 0.55, 0.7, 0.25,
                    font_size=8, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 1 — TITLE
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide1 = new_slide()
fill_bg(slide1, PURPLE)
add_rect(slide1, 0, 5.5, 13.33, 2.0, fill=DARK2)
add_rect(slide1, 0, 0, 13.33, 0.45, fill=DARK2)

add_textbox(slide1, "Field Manager Pro", 0.5, 1.6, 12.3, 1.6,
            font_size=68, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide1, "User Guide & Role Responsibilities", 0.5, 3.3, 12.3, 0.7,
            font_size=26, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide1, "fieldmanagerpro.app", 0.5, 6.15, 12.3, 0.5,
            font_size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 2 — WHAT IS FIELD MANAGER PRO?
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide2 = new_slide()
fill_bg(slide2, DARK_BG)
heading(slide2, "What is Field Manager Pro?")
divider(slide2, 1.15)
add_textbox(slide2,
    "A mobile-first platform for managing field teams, schedules, tasks, and time tracking — all in one place.",
    0.4, 1.3, 7.8, 0.7, font_size=16, color=LIGHT_GRAY)

features = [
    "Clock in/out tracking",
    "Staff scheduling",
    "Task management",
    "Time card review",
    "Flags & alerts",
    "Store inspections",
    "Expense tracking",
]
bullet_list(slide2, features, l=0.5, t=2.1, w=6.5, font_size=15, spacing=0.44)

# Note box
note = add_rect(slide2, 0.4, 6.4, 7.5, 0.65, fill=MED_GRAY, radius=True)
add_textbox(slide2, "📌  This guide is organized by role. Jump to your section.",
            0.55, 6.45, 7.3, 0.5, font_size=13, color=WHITE)

# Right panel — feature icon grid
icons = ["🕐 Clock In/Out", "📅 Schedule", "✅ Tasks", "📝 Timecards",
         "🚩 Flags", "🏪 Inspections", "💳 Expenses"]
cols = 2
col_w = 2.4
row_h = 0.72
for idx, icon in enumerate(icons):
    col = idx % cols
    row = idx // cols
    bx = 8.2 + col * col_w
    by = 1.35 + row * row_h
    b = add_rect(slide2, bx, by, col_w - 0.15, row_h - 0.1, fill=DARK2, radius=True)
    add_textbox(slide2, icon, bx + 0.1, by + 0.08, col_w - 0.3, row_h - 0.15,
                font_size=13, color=WHITE)

add_slide_number(slide2, 2)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 3 — APP OVERVIEW / NAVIGATION
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide3 = new_slide()
fill_bg(slide3, DARK_BG)
heading(slide3, "App Overview & Navigation")
divider(slide3, 1.15)

steps = [
    ("1.", "Go to fieldmanagerpro.app on your phone or computer"),
    ("2.", "Log in with your username and password"),
    ("3.", "Use the navigation bar at the BOTTOM of the screen to switch sections"),
    ("4.", "Tap icons to jump between Home, Schedule, Tasks, Team, and more"),
]
for i, (num, txt) in enumerate(steps):
    y = 1.4 + i * 0.66
    add_rect(slide3, 0.4, y, 0.38, 0.45, fill=PURPLE, radius=True)
    add_textbox(slide3, num, 0.42, y + 0.03, 0.35, 0.38,
                font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide3, txt, 0.9, y + 0.04, 6.2, 0.45, font_size=14, color=WHITE)

add_textbox(slide3, "💡  If you forget your password, tap 'Forgot your password?' on the login screen.",
            0.4, 4.2, 6.8, 0.5, font_size=13, color=LIGHT_GRAY, italic=True)

# Phone mockup with nav bar
pl, pt, pw, ph = phone_frame(slide3, l=7.6, t=0.55, w=4.9, h=6.0)

# Status bar text
add_textbox(slide3, "9:41 AM", pl + 0.15, pt + 0.1, 1.0, 0.25,
            font_size=9, color=LIGHT_GRAY)
add_textbox(slide3, "●●●●●", pl + pw - 0.9, pt + 0.1, 0.8, 0.25,
            font_size=9, color=LIGHT_GRAY)

# App header bar
add_rect(slide3, pl, pt + 0.3, pw, 0.55, fill=DARK2)
add_textbox(slide3, "Field Manager Pro", pl + 0.15, pt + 0.35, pw - 0.3, 0.45,
            font_size=14, bold=True, color=WHITE)

# Home content area - simple summary cards
card_items = ["Welcome back!", "3 Shifts This Week", "2 Tasks Pending", "No Flags"]
card_colors = [DARK2, MED_GRAY, MED_GRAY, DARK2]
for ci, (ctext, ccol) in enumerate(zip(card_items, card_colors)):
    cy = pt + 1.0 + ci * 0.82
    add_rect(slide3, pl + 0.15, cy, pw - 0.3, 0.65, fill=ccol, radius=True)
    add_textbox(slide3, ctext, pl + 0.3, cy + 0.12, pw - 0.6, 0.42,
                font_size=11, color=WHITE)

# Nav bar at bottom
nav_labels = ["Home", "Schedule", "Tasks", "Team", "More"]
nav_bar(slide3, pl, pt + ph - 0.95, pw, 0.95, nav_labels)

add_slide_number(slide3, 3)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 4 — GENERAL RULES
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide4 = new_slide()
fill_bg(slide4, DARK_BG)
heading(slide4, "General Rules — All Roles")
divider(slide4, 1.15)

# Red warning box
warn = add_rect(slide4, 0.4, 1.35, 12.5, 1.55, fill=RGBColor(0x7f, 0x1d, 0x1d), radius=True)
add_rect(slide4, 0.4, 1.35, 0.12, 1.55, fill=RED)
add_textbox(slide4, "⚠  TIME FRAUD POLICY", 0.65, 1.42, 11.8, 0.45,
            font_size=16, bold=True, color=RED)
add_textbox(slide4,
    "Any attempt to share login credentials or clock in another employee will result in IMMEDIATE TERMINATION.",
    0.65, 1.82, 11.8, 0.65, font_size=14, color=WHITE)

rules = [
    ("🚫", "Do not share your login credentials with anyone."),
    ("📧", "App Issues: If something isn't working, email your Owner with a description and screenshots of the error. The Owner will escalate to the developer."),
    ("📍", "GPS location is recorded at every clock-in and clock-out. Location must match your work location."),
    ("🔒", "Your account is personal. You are responsible for all activity under your login."),
]
for i, (icon, rule) in enumerate(rules):
    y = 3.1 + i * 0.95
    add_rect(slide4, 0.4, y, 0.5, 0.7, fill=DARK2, radius=True)
    add_textbox(slide4, icon, 0.42, y + 0.08, 0.45, 0.5,
                font_size=16, align=PP_ALIGN.CENTER)
    add_textbox(slide4, rule, 1.05, y + 0.08, 11.5, 0.65, font_size=13, color=LIGHT_GRAY)

add_slide_number(slide4, 4)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SECTION 2: EMPLOYEE
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

# Slide 5 — Employee Section Header
section_header_slide("Employee", "Your guide to using Field Manager Pro")

# ─── Slide 6 — Employee: Getting Started ──────────────────────────────────────
slide6 = new_slide()
fill_bg(slide6, DARK_BG)
heading(slide6, "Employee: Getting Started")
divider(slide6, 1.15)

steps6 = [
    "Open your browser or the app on your phone.",
    "Go to fieldmanagerpro.app",
    "Enter your Username and Password (provided by your manager).",
    "Tap 'Sign In' to access your dashboard.",
]
bullet_list(slide6, steps6, l=0.5, t=1.35, w=6.6, font_size=15, spacing=0.6)

# Login mockup
pl, pt, pw, ph = phone_frame(slide6, l=7.55, t=0.55, w=4.95, h=6.0)
add_textbox(slide6, "Field Manager Pro", pl + 0.3, pt + 0.55, pw - 0.6, 0.55,
            font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Username field
add_rect(slide6, pl + 0.25, pt + 1.35, pw - 0.5, 0.55, fill=MED_GRAY, radius=True)
add_textbox(slide6, "Username", pl + 0.4, pt + 1.42, pw - 0.8, 0.4,
            font_size=12, color=LIGHT_GRAY)

# Password field
add_rect(slide6, pl + 0.25, pt + 2.1, pw - 0.5, 0.55, fill=MED_GRAY, radius=True)
add_textbox(slide6, "Password", pl + 0.4, pt + 2.17, pw - 0.8, 0.4,
            font_size=12, color=LIGHT_GRAY)

# Sign In button
btn = add_rect(slide6, pl + 0.25, pt + 2.9, pw - 0.5, 0.6, fill=PURPLE, radius=True)
add_textbox(slide6, "Sign In", pl + 0.3, pt + 2.97, pw - 0.6, 0.45,
            font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Forgot password
add_textbox(slide6, "Forgot your password?", pl + 0.25, pt + 3.65, pw - 0.5, 0.35,
            font_size=10, color=PURPLE, align=PP_ALIGN.CENTER)

# fieldmanagerpro.app label
add_textbox(slide6, "fieldmanagerpro.app", pl + 0.2, pt + 4.8, pw - 0.4, 0.3,
            font_size=9, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_slide_number(slide6, 6)

# ─── Slide 7 — Employee: Clocking In & Out ────────────────────────────────────
slide7 = new_slide()
fill_bg(slide7, DARK_BG)
heading(slide7, "Employee: Clocking In & Out")
divider(slide7, 1.15)

steps7 = [
    "Open the app and tap the Clock icon in the nav bar.",
    "Make sure you are INSIDE the store and ready to work.",
    "Tap the green 'Clock In' button.",
    "To clock out: tap 'Clock Out' BEFORE leaving the store.",
]
bullet_list(slide7, steps7, l=0.5, t=1.35, w=6.5, font_size=15, spacing=0.55)

# Important box
imp = add_rect(slide7, 0.4, 3.6, 6.7, 1.3, fill=RGBColor(0x05, 0x46, 0x1f), radius=True)
add_rect(slide7, 0.4, 3.6, 0.1, 1.3, fill=GREEN)
add_textbox(slide7, "IMPORTANT", 0.65, 3.67, 6.3, 0.38,
            font_size=13, bold=True, color=GREEN)
add_textbox(slide7,
    "Clock in INSIDE the store, AFTER you are ready to work.\nClock out BEFORE leaving the store.",
    0.65, 4.02, 6.3, 0.75, font_size=13, color=WHITE)

add_textbox(slide7, "📍  GPS location is recorded at clock-in and clock-out.",
            0.4, 5.1, 6.7, 0.4, font_size=12, color=LIGHT_GRAY, italic=True)

# Phone mockup — Clock screen
pl, pt, pw, ph = phone_frame(slide7, l=7.55, t=0.55, w=4.95, h=6.0)
add_textbox(slide7, "Clock", pl + 0.3, pt + 0.45, pw - 0.6, 0.5,
            font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide7, "9:41 AM", pl + 0.3, pt + 1.2, pw - 0.6, 0.6,
            font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide7, "Tuesday, April 15", pl + 0.3, pt + 1.82, pw - 0.6, 0.35,
            font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Big clock-in button
btn7 = add_rect(slide7, pl + 0.5, pt + 2.4, pw - 1.0, 1.0, fill=GREEN, radius=True)
add_textbox(slide7, "Clock In", pl + 0.5, pt + 2.6, pw - 1.0, 0.6,
            font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(slide7, "Not clocked in", pl + 0.3, pt + 3.7, pw - 0.6, 0.4,
            font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

nav_bar(slide7, pl, pt + ph - 0.95, pw, 0.95, ["Home", "Clock", "Tasks", "History"])

add_slide_number(slide7, 7)

# ─── Slide 8 — Employee: Viewing Your Schedule ───────────────────────────────
slide8 = new_slide()
fill_bg(slide8, DARK_BG)
heading(slide8, "Employee: Viewing Your Schedule")
divider(slide8, 1.15)

bullet_list(slide8, [
    "Tap 'Schedule' in the bottom navigation bar.",
    "Your assigned shifts for the week will appear here.",
    "Each shift block shows the store name, date, and time.",
    "Contact your DM if you see a scheduling error.",
], l=0.5, t=1.35, w=6.5, font_size=15, spacing=0.6)

# Phone mockup — schedule
pl, pt, pw, ph = phone_frame(slide8, l=7.55, t=0.55, w=4.95, h=6.0)
add_textbox(slide8, "My Schedule", pl + 0.2, pt + 0.45, pw - 0.4, 0.45,
            font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

days_abbr = ["M", "T", "W", "T", "F", "S", "S"]
day_w = (pw - 0.3) / 7
for di, d in enumerate(days_abbr):
    dx = pl + 0.15 + di * day_w
    add_rect(slide8, dx, pt + 1.0, day_w - 0.05, 0.38,
             fill=MED_GRAY if di not in [0, 2, 4] else PURPLE, radius=True)
    add_textbox(slide8, d, dx, pt + 1.06, day_w - 0.05, 0.28,
                font_size=10, bold=(di in [0,2,4]), color=WHITE, align=PP_ALIGN.CENTER)

# Shift cards
shift_days = [(0, "Mon Apr 14"), (2, "Wed Apr 16"), (4, "Fri Apr 18")]
for idx, (di, label) in enumerate(shift_days):
    sy = pt + 1.65 + idx * 1.25
    sc = add_rect(slide8, pl + 0.15, sy, pw - 0.3, 1.0, fill=DARK2, radius=True)
    add_rect(slide8, pl + 0.15, sy, 0.1, 1.0, fill=PURPLE)
    add_textbox(slide8, "Store Shift", pl + 0.35, sy + 0.08, pw - 0.55, 0.35,
                font_size=12, bold=True, color=WHITE)
    add_textbox(slide8, label, pl + 0.35, sy + 0.42, pw - 0.55, 0.28,
                font_size=10, color=LIGHT_GRAY)
    add_textbox(slide8, "9:00 AM – 5:00 PM", pl + 0.35, sy + 0.65, pw - 0.55, 0.28,
                font_size=10, color=LIGHT_GRAY)

nav_bar(slide8, pl, pt + ph - 0.95, pw, 0.95, ["Home", "Schedule", "Tasks", "Team"])

add_slide_number(slide8, 8)

# ─── Slide 9 — Employee: Checklist ───────────────────────────────────────────
slide9 = new_slide()
fill_bg(slide9, DARK_BG)
heading(slide9, "Employee: Checklist")
divider(slide9, 1.15)

bullet_list(slide9, [
    "Tap 'Tasks' or 'Checklist' in the bottom navigation bar.",
    "Your daily tasks assigned by your DM will appear here.",
    "Tap each item to mark it as complete.",
    "Add a photo or note if required by your manager.",
], l=0.5, t=1.35, w=6.5, font_size=15, spacing=0.6)

# Phone mockup — checklist
pl, pt, pw, ph = phone_frame(slide9, l=7.55, t=0.55, w=4.95, h=6.0)
add_textbox(slide9, "Daily Checklist", pl + 0.2, pt + 0.45, pw - 0.4, 0.45,
            font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

checklist_items = [
    ("Stock shelves in aisle 3", True),
    ("Complete opening report", True),
    ("Check cooler temperatures", False),
    ("Submit end-of-day photos", False),
]
for ci, (citem, done) in enumerate(checklist_items):
    cy = pt + 1.1 + ci * 1.05
    add_rect(slide9, pl + 0.15, cy, pw - 0.3, 0.85, fill=DARK2, radius=True)
    # circle checkbox
    cbox = slide9.shapes.add_shape(
        9, Inches(pl + 0.27), Inches(cy + 0.2), Inches(0.42), Inches(0.42)
    )
    cbox.fill.solid()
    cbox.fill.fore_color.rgb = GREEN if done else MED_GRAY
    cbox.line.fill.background()
    if done:
        add_textbox(slide9, "✓", pl + 0.27, cy + 0.18, 0.42, 0.38,
                    font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide9, citem, pl + 0.8, cy + 0.22, pw - 1.05, 0.42,
                font_size=11, color=WHITE if not done else LIGHT_GRAY)

nav_bar(slide9, pl, pt + ph - 0.95, pw, 0.95, ["Home", "Schedule", "Tasks", "History"])

add_slide_number(slide9, 9)

# ─── Slide 10 — Employee: Time History ───────────────────────────────────────
slide10 = new_slide()
fill_bg(slide10, DARK_BG)
heading(slide10, "Employee: Time History")
divider(slide10, 1.15)

bullet_list(slide10, [
    "Tap 'History' or 'Time History' in the navigation bar.",
    "View a log of all your past clock-in and clock-out records.",
    "Each entry shows the date, time in, time out, and total hours.",
    "If you believe there is an error in your record, notify your DM immediately.",
], l=0.5, t=1.35, w=6.5, font_size=15, spacing=0.65)

# Phone mockup — time history list
pl, pt, pw, ph = phone_frame(slide10, l=7.55, t=0.55, w=4.95, h=6.0)
add_textbox(slide10, "Time History", pl + 0.2, pt + 0.45, pw - 0.4, 0.45,
            font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

history = [
    ("Apr 15, Tue", "9:02 AM", "5:04 PM", "8h 2m"),
    ("Apr 14, Mon", "8:58 AM", "5:01 PM", "8h 3m"),
    ("Apr 12, Sat", "10:00 AM", "3:00 PM", "5h 0m"),
    ("Apr 11, Fri",  "9:05 AM", "5:10 PM", "8h 5m"),
]
for hi, (date, tin, tout, total) in enumerate(history):
    hy = pt + 1.1 + hi * 1.1
    add_rect(slide10, pl + 0.15, hy, pw - 0.3, 0.9, fill=DARK2, radius=True)
    add_textbox(slide10, date, pl + 0.28, hy + 0.06, pw - 0.55, 0.32,
                font_size=11, bold=True, color=WHITE)
    add_textbox(slide10, f"In: {tin}  |  Out: {tout}", pl + 0.28, hy + 0.38, pw - 0.55, 0.28,
                font_size=10, color=LIGHT_GRAY)
    # Total pill
    pill = add_rect(slide10, pl + pw - 1.3, hy + 0.28, 0.9, 0.32, fill=PURPLE, radius=True)
    add_textbox(slide10, total, pl + pw - 1.3, hy + 0.3, 0.9, 0.28,
                font_size=9, color=WHITE, align=PP_ALIGN.CENTER)

nav_bar(slide10, pl, pt + ph - 0.95, pw, 0.95, ["Home", "Clock", "Tasks", "History"])

add_slide_number(slide10, 10)

# ─── Slide 11 — Employee: Expenses ───────────────────────────────────────────
slide11 = new_slide()
fill_bg(slide11, DARK_BG)
heading(slide11, "Employee: Expenses")
divider(slide11, 1.15)

bullet_list(slide11, [
    "Tap 'Expenses' in the navigation bar.",
    "Fill in: Date, Category, Amount, Description.",
    "Attach a photo of your receipt.",
    "Tap 'Submit' — your manager will review and approve or reject.",
    "You will be notified of the decision via the app.",
], l=0.5, t=1.35, w=6.5, font_size=15, spacing=0.58)

# Phone mockup — expense form
pl, pt, pw, ph = phone_frame(slide11, l=7.55, t=0.55, w=4.95, h=6.0)
add_textbox(slide11, "New Expense", pl + 0.2, pt + 0.45, pw - 0.4, 0.45,
            font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

fields11 = [
    ("Date", "Apr 15, 2026"),
    ("Category", "Supplies  ▾"),
    ("Amount", "$  0.00"),
    ("Description", "Enter details..."),
]
for fi, (flbl, fval) in enumerate(fields11):
    fy = pt + 1.1 + fi * 0.88
    add_textbox(slide11, flbl, pl + 0.22, fy, pw - 0.45, 0.28,
                font_size=9, color=LIGHT_GRAY)
    add_rect(slide11, pl + 0.22, fy + 0.27, pw - 0.44, 0.5, fill=MED_GRAY, radius=True)
    add_textbox(slide11, fval, pl + 0.36, fy + 0.32, pw - 0.72, 0.38,
                font_size=11, color=WHITE)

# Photo upload box
pu = add_rect(slide11, pl + 0.22, pt + 4.7, pw - 0.44, 0.75, fill=MED_GRAY, radius=True)
add_textbox(slide11, "📷  Attach Receipt Photo", pl + 0.22, pt + 4.84, pw - 0.44, 0.45,
            font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Submit button
add_rect(slide11, pl + 0.22, pt + 5.55, pw - 0.44, 0.5, fill=PURPLE, radius=True)
add_textbox(slide11, "Submit Expense", pl + 0.22, pt + 5.62, pw - 0.44, 0.38,
            font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_slide_number(slide11, 11)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SECTION 3: DM
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

# Slide 12 — DM Section Header
section_header_slide("DM (District Manager)", "Your guide to using Field Manager Pro")

# ─── Slide 13 — DM: Overview & Responsibilities ──────────────────────────────
slide13 = new_slide()
fill_bg(slide13, DARK_BG)
heading(slide13, "DM: Overview & Responsibilities")
divider(slide13, 1.15)

dm_resp = [
    "Manage your assigned team of employees",
    "Post store schedule 2 WEEKS in advance at all times",
    "Your personal schedule must be posted 1 WEEK in advance — or you cannot clock in",
    "Review and validate timecards every Monday",
    "Check the app for new tasks and flags DAILY",
    "Clock in when you start; clock out ONLY when your shift is fully complete",
    "You are responsible for your team's performance and schedule compliance",
]
bullet_list(slide13, dm_resp, l=0.5, t=1.3, w=12.2, font_size=14, spacing=0.66)

add_slide_number(slide13, 13)

# ─── Slide 14 — DM: Staff Scheduling Part 1 ──────────────────────────────────
slide14 = new_slide()
fill_bg(slide14, DARK_BG)
heading(slide14, "DM: Staff Scheduling — How To")
divider(slide14, 1.15)

steps14 = [
    ("1.", "Navigate to 'Staff Schedule' in the nav bar."),
    ("2.", "Tap the '+' button to add a shift for an employee."),
    ("3.", "Select the employee, date, start time, and end time."),
    ("4.", "Tap 'Publish' to make the schedule visible to your team."),
]
for i, (num, txt) in enumerate(steps14):
    y = 1.35 + i * 0.7
    add_rect(slide14, 0.4, y, 0.38, 0.48, fill=PURPLE, radius=True)
    add_textbox(slide14, num, 0.42, y + 0.04, 0.35, 0.38,
                font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide14, txt, 0.9, y + 0.07, 6.3, 0.48, font_size=14, color=WHITE)

# Schedule grid mockup
gx, gy, gw, gh = 7.3, 1.0, 5.7, 6.1
add_rect(slide14, gx, gy, gw, gh, fill=DARK2, radius=True)
add_textbox(slide14, "Staff Schedule — Week of Apr 14", gx + 0.15, gy + 0.1, gw - 0.3, 0.4,
            font_size=12, bold=True, color=WHITE)

days14 = ["Mon", "Tue", "Wed", "Thu", "Fri", "Sat", "Sun"]
cell_w = (gw - 1.2) / 7
header_y = gy + 0.6
for di, d in enumerate(days14):
    dx = gx + 1.2 + di * cell_w
    add_rect(slide14, dx, header_y, cell_w - 0.05, 0.4, fill=PURPLE, radius=False)
    add_textbox(slide14, d, dx, header_y + 0.05, cell_w - 0.05, 0.3,
                font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

names14 = ["J. Smith", "A. Jones", "M. Lee", "P. Davis"]
cell_h = 0.9
for ni, name in enumerate(names14):
    ny = header_y + 0.4 + ni * cell_h
    add_rect(slide14, gx + 0.05, ny, 1.1, cell_h - 0.05, fill=MED_GRAY)
    add_textbox(slide14, name, gx + 0.08, ny + 0.28, 1.0, 0.35, font_size=9, color=WHITE)
    # fill some cells
    filled_days = {0:[0,2,4], 1:[1,3,5], 2:[0,1,4,6], 3:[2,4]}
    for di in filled_days.get(ni, []):
        dx = gx + 1.2 + di * cell_w
        add_rect(slide14, dx + 0.02, ny + 0.08, cell_w - 0.09, cell_h - 0.22,
                 fill=PURPLE, radius=True)
        add_textbox(slide14, "9-5", dx + 0.02, ny + 0.18, cell_w - 0.09, 0.45,
                    font_size=8, color=WHITE, align=PP_ALIGN.CENTER)

add_slide_number(slide14, 14)

# ─── Slide 15 — DM: Scheduling Requirements ──────────────────────────────────
slide15 = new_slide()
fill_bg(slide15, DARK_BG)
heading(slide15, "DM: Scheduling Requirements")
divider(slide15, 1.15)

# Req 1 box
req1 = add_rect(slide15, 0.4, 1.3, 5.8, 2.0, fill=DARK2, radius=True)
add_rect(slide15, 0.4, 1.3, 0.1, 2.0, fill=PURPLE)
add_textbox(slide15, "REQUIREMENT 1", 0.65, 1.38, 5.4, 0.38,
            font_size=12, bold=True, color=PURPLE)
add_textbox(slide15, "Employee schedules must be posted\n2 weeks in advance.", 0.65, 1.75, 5.4, 0.7,
            font_size=14, bold=True, color=WHITE)
add_textbox(slide15, "Example: If today is Mon Apr 14, your schedule\nmust be complete through at least Sun Apr 27.",
            0.65, 2.42, 5.4, 0.75, font_size=12, color=LIGHT_GRAY, italic=True)

# Req 2 box
req2 = add_rect(slide15, 6.6, 1.3, 6.3, 2.0, fill=DARK2, radius=True)
add_rect(slide15, 6.6, 1.3, 0.1, 2.0, fill=ORANGE)
add_textbox(slide15, "REQUIREMENT 2", 6.85, 1.38, 5.8, 0.38,
            font_size=12, bold=True, color=ORANGE)
add_textbox(slide15, "Your personal schedule must be posted\nat least 1 week in advance.", 6.85, 1.75, 5.8, 0.7,
            font_size=14, bold=True, color=WHITE)
add_textbox(slide15, "Example: If today is Mon Apr 14, your schedule\nmust be set through at least Sun Apr 20.",
            6.85, 2.42, 5.8, 0.75, font_size=12, color=LIGHT_GRAY, italic=True)

# Consequence box
cons = add_rect(slide15, 0.4, 3.55, 12.5, 0.85, fill=RGBColor(0x7f, 0x1d, 0x1d), radius=True)
add_textbox(slide15, "⚠  CONSEQUENCE:  If your personal schedule is not loaded, you will be UNABLE to clock in.",
            0.6, 3.67, 12.2, 0.6, font_size=14, bold=True, color=RED)

# Two-week calendar mockup
cal_x, cal_y = 0.4, 4.6
cal_w, cal_h = 12.5, 2.55
add_rect(slide15, cal_x, cal_y, cal_w, cal_h, fill=DARK2, radius=True)
add_textbox(slide15, "April 2026", cal_x + 0.15, cal_y + 0.08, cal_w - 0.3, 0.38,
            font_size=12, bold=True, color=WHITE)

days_cal = ["Mon", "Tue", "Wed", "Thu", "Fri", "Sat", "Sun"]
cal_cw = (cal_w - 0.3) / 7
for di, d in enumerate(days_cal):
    dx = cal_x + 0.15 + di * cal_cw
    add_textbox(slide15, d, dx, cal_y + 0.48, cal_cw, 0.28,
                font_size=9, bold=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Two weeks of dates (Apr 14-27)
dates_row = [14, 15, 16, 17, 18, 19, 20]
dates_row2 = [21, 22, 23, 24, 25, 26, 27]
for di, (d1, d2) in enumerate(zip(dates_row, dates_row2)):
    dx = cal_x + 0.15 + di * cal_cw
    # Week 1 = orange (personal), Week 2 = purple (team)
    c1 = add_rect(slide15, dx + 0.05, cal_y + 0.82, cal_cw - 0.1, 0.55, fill=ORANGE, radius=True)
    add_textbox(slide15, str(d1), dx + 0.05, cal_y + 0.89, cal_cw - 0.1, 0.38,
                font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    c2 = add_rect(slide15, dx + 0.05, cal_y + 1.45, cal_cw - 0.1, 0.55, fill=PURPLE, radius=True)
    add_textbox(slide15, str(d2), dx + 0.05, cal_y + 1.52, cal_cw - 0.1, 0.38,
                font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(slide15, "Orange = Personal (1 wk req.)", 0.5, 7.15, 4.0, 0.28,
            font_size=9, color=ORANGE)
add_textbox(slide15, "Purple = Team/Employees (2 wk req.)", 4.6, 7.15, 4.5, 0.28,
            font_size=9, color=PURPLE)

add_slide_number(slide15, 15)

# ─── Slide 16 — DM: Timecard Review ──────────────────────────────────────────
slide16 = new_slide()
fill_bg(slide16, DARK_BG)
heading(slide16, "DM: Timecard Review")
divider(slide16, 1.15)

bullet_list(slide16, [
    "Navigate to 'Timecards' in the nav bar.",
    "Review each employee's clock-in and clock-out times weekly.",
    "REQUIRED: Complete timecard review every Monday.",
    "If an entry looks incorrect: tap it → edit clock-in/out → add a note explaining the change.",
], l=0.5, t=1.35, w=6.5, font_size=14, spacing=0.72)

# Phone mockup — timecard list
pl, pt, pw, ph = phone_frame(slide16, l=7.55, t=0.55, w=4.95, h=6.0)
add_textbox(slide16, "Timecards", pl + 0.2, pt + 0.45, pw - 0.4, 0.45,
            font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

tc_entries = [
    ("J. Smith", "Apr 15", "9:02 AM", "5:04 PM"),
    ("A. Jones", "Apr 15", "8:55 AM", "5:10 PM"),
    ("M. Lee",   "Apr 15", "9:15 AM", "—"),
]
for ti, (name, date, tin, tout) in enumerate(tc_entries):
    ty = pt + 1.05 + ti * 1.4
    add_rect(slide16, pl + 0.15, ty, pw - 0.3, 1.2, fill=DARK2, radius=True)
    add_textbox(slide16, name, pl + 0.28, ty + 0.08, pw - 0.7, 0.32,
                font_size=12, bold=True, color=WHITE)
    add_textbox(slide16, date, pl + 0.28, ty + 0.38, pw - 0.7, 0.28,
                font_size=10, color=LIGHT_GRAY)
    add_textbox(slide16, f"In: {tin}", pl + 0.28, ty + 0.65, 1.6, 0.28,
                font_size=10, color=GREEN if tin != "—" else RED)
    add_textbox(slide16, f"Out: {tout}", pl + 0.28, ty + 0.9, 1.6, 0.28,
                font_size=10, color=GREEN if tout != "—" else RED)
    # Edit button
    edit_btn = add_rect(slide16, pl + pw - 1.1, ty + 0.42, 0.8, 0.38, fill=PURPLE, radius=True)
    add_textbox(slide16, "Edit", pl + pw - 1.1, ty + 0.46, 0.8, 0.3,
                font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

nav_bar(slide16, pl, pt + ph - 0.95, pw, 0.95, ["Home", "Schedule", "Timecards", "Flags"])

add_slide_number(slide16, 16)

# ─── Slide 17 — DM: Tasks ────────────────────────────────────────────────────
slide17 = new_slide()
fill_bg(slide17, DARK_BG)
heading(slide17, "DM: Tasks")
divider(slide17, 1.15)

bullet_list(slide17, [
    "Navigate to 'Tasks' in the nav bar.",
    "Tasks are assigned by your Owner, Sales Director, or RDM.",
    "Check for new tasks DAILY.",
    "To complete: tap the circle → add a note (optional) → add a photo (optional) → tap 'Mark Complete'.",
], l=0.5, t=1.35, w=6.5, font_size=14, spacing=0.72)

# Phone mockup — task card
pl, pt, pw, ph = phone_frame(slide17, l=7.55, t=0.55, w=4.95, h=6.0)
add_textbox(slide17, "Tasks", pl + 0.2, pt + 0.45, pw - 0.4, 0.45,
            font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

task_items = [
    ("Submit Weekly Report", "Owner", "Apr 18", ORANGE, False),
    ("Train new staff member", "Sales Dir.", "Apr 20", PURPLE, False),
    ("Inventory count complete", "Owner", "Apr 14", GREEN, True),
]
for tki, (ttitle, tassigner, tdue, tcol, tdone) in enumerate(task_items):
    ty = pt + 1.05 + tki * 1.6
    add_rect(slide17, pl + 0.15, ty, pw - 0.3, 1.4, fill=DARK2, radius=True)
    # Checkbox circle
    ccirc = slide17.shapes.add_shape(
        9, Inches(pl + 0.28), Inches(ty + 0.48), Inches(0.42), Inches(0.42)
    )
    ccirc.fill.solid()
    ccirc.fill.fore_color.rgb = GREEN if tdone else MED_GRAY
    ccirc.line.fill.background()
    if tdone:
        add_textbox(slide17, "✓", pl + 0.28, ty + 0.46, 0.42, 0.38,
                    font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide17, ttitle, pl + 0.82, ty + 0.08, pw - 1.05, 0.4,
                font_size=12, bold=True, color=WHITE)
    add_textbox(slide17, f"From: {tassigner}", pl + 0.82, ty + 0.48, 1.5, 0.28,
                font_size=9, color=LIGHT_GRAY)
    # Due badge
    due_b = add_rect(slide17, pl + pw - 1.35, ty + 0.08, 1.1, 0.32, fill=tcol, radius=True)
    add_textbox(slide17, f"Due {tdue}", pl + pw - 1.35, ty + 0.1, 1.1, 0.28,
                font_size=8, color=WHITE, align=PP_ALIGN.CENTER)

    if not tdone:
        # Mark Complete button
        mc = add_rect(slide17, pl + 0.82, ty + 0.92, pw - 1.05, 0.38, fill=PURPLE, radius=True)
        add_textbox(slide17, "Mark Complete", pl + 0.82, ty + 0.96, pw - 1.05, 0.3,
                    font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

nav_bar(slide17, pl, pt + ph - 0.95, pw, 0.95, ["Home", "Schedule", "Tasks", "Flags"])

add_slide_number(slide17, 17)

# ─── Slide 18 — DM: Flags ────────────────────────────────────────────────────
slide18 = new_slide()
fill_bg(slide18, DARK_BG)
heading(slide18, "DM: Flags")
divider(slide18, 1.15)

bullet_list(slide18, [
    "Navigate to 'Flags' in the nav bar.",
    "Flags are auto-created when unusual time activity is detected.",
    "Examples: missed clock-out, overtime, late clock-in.",
    "Review flags DAILY and resolve them with a note.",
], l=0.5, t=1.35, w=6.5, font_size=14, spacing=0.72)

flag_types = [
    ("Missed Clock-Out", RED),
    ("Late Clock-In (>15 min)", ORANGE),
    ("Overtime Alert", ORANGE),
]
for fi, (ftype, fcol) in enumerate(flag_types):
    fy = 3.6 + fi * 0.65
    add_rect(slide18, 0.4, fy, 0.12, 0.48, fill=fcol)
    add_textbox(slide18, ftype, 0.65, fy + 0.08, 5.8, 0.35, font_size=13, color=WHITE)

# Phone mockup — flag card
pl, pt, pw, ph = phone_frame(slide18, l=7.55, t=0.55, w=4.95, h=6.0)
add_textbox(slide18, "Flags", pl + 0.2, pt + 0.45, pw - 0.4, 0.45,
            font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

flags = [
    ("J. Smith", "Missed Clock-Out", "Apr 15", RED),
    ("A. Jones", "Late Clock-In", "Apr 14", ORANGE),
]
for fli, (fname, ftype, fdate, fcol) in enumerate(flags):
    fly = pt + 1.05 + fli * 2.3
    add_rect(slide18, pl + 0.15, fly, pw - 0.3, 2.05, fill=DARK2, radius=True)
    add_rect(slide18, pl + 0.15, fly, 0.1, 2.05, fill=fcol)
    # Flag type badge
    fb = add_rect(slide18, pl + 0.35, fly + 0.1, pw - 0.75, 0.38, fill=fcol, radius=True)
    add_textbox(slide18, ftype, pl + 0.35, fly + 0.12, pw - 0.75, 0.3,
                font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide18, fname, pl + 0.35, fly + 0.58, pw - 0.75, 0.35,
                font_size=13, bold=True, color=WHITE)
    add_textbox(slide18, fdate, pl + 0.35, fly + 0.9, pw - 0.75, 0.3,
                font_size=10, color=LIGHT_GRAY)
    # Resolve button
    res_btn = add_rect(slide18, pl + 0.35, fly + 1.48, pw - 0.75, 0.42, fill=PURPLE, radius=True)
    add_textbox(slide18, "Resolve", pl + 0.35, fly + 1.52, pw - 0.75, 0.34,
                font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

nav_bar(slide18, pl, pt + ph - 0.95, pw, 0.95, ["Home", "Schedule", "Tasks", "Flags"])

add_slide_number(slide18, 18)

# ─── Slide 19 — DM: Clocking In & Out ────────────────────────────────────────
slide19 = new_slide()
fill_bg(slide19, DARK_BG)
heading(slide19, "DM: Clocking In & Out")
divider(slide19, 1.15)

bullet_list(slide19, [
    "Clock in when you START your workday.",
    "Clock out ONLY when your shift is completely finished for the day.",
    "GPS location is recorded — must match your work location.",
    "Do not clock out during breaks unless your shift is fully complete.",
    "If your personal schedule is not posted, you will be unable to clock in.",
], l=0.5, t=1.35, w=6.5, font_size=14, spacing=0.65)

# Phone mockup — Clock screen (same as employee, DM context)
pl, pt, pw, ph = phone_frame(slide19, l=7.55, t=0.55, w=4.95, h=6.0)
add_textbox(slide19, "Clock", pl + 0.3, pt + 0.45, pw - 0.6, 0.5,
            font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide19, "9:05 AM", pl + 0.3, pt + 1.2, pw - 0.6, 0.6,
            font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide19, "Tuesday, April 15", pl + 0.3, pt + 1.82, pw - 0.6, 0.35,
            font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

btn19 = add_rect(slide19, pl + 0.5, pt + 2.4, pw - 1.0, 1.0, fill=GREEN, radius=True)
add_textbox(slide19, "Clock In", pl + 0.5, pt + 2.6, pw - 1.0, 0.6,
            font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(slide19, "📍 GPS: Main Street Store", pl + 0.3, pt + 3.7, pw - 0.6, 0.4,
            font_size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

nav_bar(slide19, pl, pt + ph - 0.95, pw, 0.95, ["Home", "Clock", "Tasks", "Flags"])

add_slide_number(slide19, 19)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SECTION 4: OPS MANAGER
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

# Slide 20 — Ops Manager Section Header
section_header_slide("Ops Manager", "Field Manager Pro")

# ─── Slide 21 — Ops Manager: Overview ────────────────────────────────────────
slide21 = new_slide()
fill_bg(slide21, DARK_BG)
heading(slide21, "Ops Manager: Overview & Responsibilities")
divider(slide21, 1.15)

bullet_list(slide21, [
    "Assigned tasks directly by the Owner.",
    "Access level is determined by what the Owner delegates.",
    "Responsibilities vary based on Owner's directives.",
    "May manage scheduling, timecards, flags, and tasks depending on assignment.",
    "Follow all DM-level expectations for any area you are delegated.",
    "Report any issues or blockers directly to the Owner.",
], l=0.5, t=1.35, w=12.2, font_size=14, spacing=0.68)

add_slide_number(slide21, 21)

# ─── Slide 22 — Ops Manager: Key Features ────────────────────────────────────
slide22 = new_slide()
fill_bg(slide22, DARK_BG)
heading(slide22, "Ops Manager: Key Features Available")
divider(slide22, 1.15)

features22 = [
    ("📅", "Staff Schedule", "View and manage team schedules"),
    ("📝", "Timecards", "View and edit employee time records"),
    ("✅", "Tasks", "Receive and complete assigned tasks"),
    ("🚩", "Flags", "View and resolve time-related flags"),
    ("👥", "Team Management", "View the team roster"),
]
for fi, (icon, fname, fdesc) in enumerate(features22):
    fy = 1.35 + fi * 1.0
    add_rect(slide22, 0.4, fy, 7.5, 0.8, fill=DARK2, radius=True)
    add_textbox(slide22, icon, 0.55, fy + 0.15, 0.5, 0.5, font_size=20)
    add_textbox(slide22, fname, 1.15, fy + 0.05, 2.5, 0.38,
                font_size=14, bold=True, color=WHITE)
    add_textbox(slide22, fdesc, 1.15, fy + 0.42, 6.5, 0.32,
                font_size=12, color=LIGHT_GRAY)

# Nav mockup
pl22, pt22, pw22, ph22 = phone_frame(slide22, l=8.5, t=1.2, w=4.5, h=5.0)
add_textbox(slide22, "Ops Manager View", pl22 + 0.1, pt22 + 0.42, pw22 - 0.2, 0.4,
            font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
nav_labels22 = ["Schedule", "Timecards", "Tasks", "Flags", "Team"]
nav_bar(slide22, pl22, pt22 + ph22 - 0.95, pw22, 0.95, nav_labels22)

add_slide_number(slide22, 22)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SECTION 5: SALES DIRECTOR
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

# Slide 23 — Sales Director Section Header
section_header_slide("Sales Director", "Field Manager Pro")

# ─── Slide 24 — Sales Director: Overview ─────────────────────────────────────
slide24 = new_slide()
fill_bg(slide24, DARK_BG)
heading(slide24, "Sales Director: Overview & Responsibilities")
divider(slide24, 1.15)

bullet_list(slide24, [
    "Assign and manage tasks for DMs across your stores.",
    "Monitor schedule compliance — all DM schedules must be 2 weeks in advance.",
    "Validate and review timecards for your DMs.",
    "Use the app as your primary management and inspection tool.",
    "Ensure task completion deadlines are met.",
    "Report unresolved flags or schedule non-compliance to the Owner.",
], l=0.5, t=1.35, w=12.2, font_size=14, spacing=0.68)

add_slide_number(slide24, 24)

# ─── Slide 25 — Sales Director: Assigning Tasks ──────────────────────────────
slide25 = new_slide()
fill_bg(slide25, DARK_BG)
heading(slide25, "Sales Director: Assigning Tasks")
divider(slide25, 1.15)

steps25 = [
    ("1.", "Navigate to 'Tasks' in the nav bar."),
    ("2.", "Tap '+ Assign Task' at the top right."),
    ("3.", "Fill in: Title, Description (optional), Assign To, Due Date."),
    ("4.", "Tap 'Assign Task' — the assignee receives an email notification."),
]
for i, (num, txt) in enumerate(steps25):
    y = 1.35 + i * 0.72
    add_rect(slide25, 0.4, y, 0.38, 0.48, fill=PURPLE, radius=True)
    add_textbox(slide25, num, 0.42, y + 0.04, 0.35, 0.38,
                font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide25, txt, 0.9, y + 0.07, 6.0, 0.48, font_size=14, color=WHITE)

# Task creation modal mockup (sheet sliding from bottom)
pl25, pt25, pw25, ph25 = phone_frame(slide25, l=7.55, t=0.55, w=4.95, h=6.0)
# Dim overlay
add_rect(slide25, pl25, pt25 + 0.3, pw25, ph25 - 0.3, fill=RGBColor(0x00, 0x00, 0x00))
# Modal sheet
modal_t = pt25 + 1.5
add_rect(slide25, pl25, modal_t, pw25, ph25 - 1.5, fill=DARK2, radius=True)
# Handle bar
add_rect(slide25, pl25 + pw25/2 - 0.3, modal_t + 0.12, 0.6, 0.08, fill=MED_GRAY, radius=True)
add_textbox(slide25, "Assign Task", pl25 + 0.2, modal_t + 0.28, pw25 - 0.4, 0.38,
            font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

modal_fields = [("Title", "Enter task title..."), ("Description", "Optional..."),
                ("Assign To", "Select DM  ▾"), ("Due Date", "Apr 18, 2026  📅")]
for mfi, (mfl, mfv) in enumerate(modal_fields):
    mfy = modal_t + 0.75 + mfi * 0.82
    add_textbox(slide25, mfl, pl25 + 0.22, mfy, pw25 - 0.44, 0.25,
                font_size=9, color=LIGHT_GRAY)
    add_rect(slide25, pl25 + 0.22, mfy + 0.23, pw25 - 0.44, 0.45, fill=MED_GRAY, radius=True)
    add_textbox(slide25, mfv, pl25 + 0.35, mfy + 0.28, pw25 - 0.70, 0.35,
                font_size=10, color=WHITE)

# Assign button
add_rect(slide25, pl25 + 0.22, modal_t + ph25 - 2.05, pw25 - 0.44, 0.55, fill=PURPLE, radius=True)
add_textbox(slide25, "Assign Task", pl25 + 0.22, modal_t + ph25 - 1.99, pw25 - 0.44, 0.43,
            font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_slide_number(slide25, 25)

# ─── Slide 26 — Sales Director: Schedule Oversight ───────────────────────────
slide26 = new_slide()
fill_bg(slide26, DARK_BG)
heading(slide26, "Sales Director: Schedule Oversight")
divider(slide26, 1.15)

bullet_list(slide26, [
    "Navigate to 'Staff Schedule' to view schedules across your stores.",
    "Verify each DM has schedules posted at least 2 weeks in advance.",
    "Navigate to 'Timecards' to review DM time entries.",
    "Follow up with DMs whose schedules are incomplete.",
    "Report persistent non-compliance to the Owner.",
], l=0.5, t=1.35, w=6.5, font_size=14, spacing=0.68)

# Schedule grid mockup (reuse grid style)
gx26, gy26, gw26, gh26 = 7.3, 1.0, 5.7, 6.1
add_rect(slide26, gx26, gy26, gw26, gh26, fill=DARK2, radius=True)
add_textbox(slide26, "Staff Schedule — All Stores", gx26 + 0.15, gy26 + 0.1, gw26 - 0.3, 0.4,
            font_size=12, bold=True, color=WHITE)

for di, d in enumerate(days14):
    dx = gx26 + 1.2 + di * cell_w
    add_rect(slide26, dx, gy26 + 0.6, cell_w - 0.05, 0.4, fill=PURPLE)
    add_textbox(slide26, d, dx, gy26 + 0.65, cell_w - 0.05, 0.3,
                font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

dm_names = ["DM: J. Brown", "DM: S. Kim", "DM: A. Patel", "DM: C. Rivera"]
for ni, name in enumerate(dm_names):
    ny = gy26 + 1.0 + ni * 1.1
    add_rect(slide26, gx26 + 0.05, ny, 1.1, 1.05, fill=MED_GRAY)
    add_textbox(slide26, name, gx26 + 0.08, ny + 0.3, 1.0, 0.45, font_size=8, color=WHITE)
    for di in [0,1,2,3,4]:
        dx = gx26 + 1.2 + di * cell_w
        col26 = PURPLE if ni < 3 else RED
        add_rect(slide26, dx + 0.02, ny + 0.15, cell_w - 0.09, 0.7, fill=col26, radius=True)
        add_textbox(slide26, "✓" if ni < 3 else "!", dx + 0.02, ny + 0.28,
                    cell_w - 0.09, 0.38, font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(slide26, "Red = Missing schedule", gx26 + 0.1, gy26 + gh26 - 0.35, 3.0, 0.3,
            font_size=9, color=RED)

add_slide_number(slide26, 26)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SECTION 6: OWNER
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

# Slide 27 — Owner Section Header
section_header_slide("Owner", "Field Manager Pro")

# ─── Slide 28 — Owner: Overview ──────────────────────────────────────────────
slide28 = new_slide()
fill_bg(slide28, DARK_BG)
heading(slide28, "Owner: Overview & Responsibilities")
divider(slide28, 1.15)

bullet_list(slide28, [
    "Primary operator of the organization in Field Manager Pro.",
    "Run company operations and inspections through the app.",
    "Assign tasks to Sales Director, Ops Manager, and DMs as needed.",
    "Oversee scheduling compliance, timecard accuracy, and flag resolution.",
    "Manage your team roster through the Team section.",
    "First point of contact for any app issues — forward to developer with screenshots.",
], l=0.5, t=1.35, w=12.2, font_size=14, spacing=0.68)

add_slide_number(slide28, 28)

# ─── Slide 29 — Owner: Team Management ───────────────────────────────────────
slide29 = new_slide()
fill_bg(slide29, DARK_BG)
heading(slide29, "Owner: Team Management")
divider(slide29, 1.15)

bullet_list(slide29, [
    "Navigate to 'Team' in the nav bar.",
    "Tap '+ Add DM' or '+ Add Employee' to create a new user.",
    "Fill in: Name, Username, Email, Temporary Password, and Role.",
    "The user receives a welcome email with their login credentials.",
    "Edit or deactivate users from the same Team screen.",
], l=0.5, t=1.35, w=6.5, font_size=14, spacing=0.7)

# Phone mockup — Team page
pl29, pt29, pw29, ph29 = phone_frame(slide29, l=7.55, t=0.55, w=4.95, h=6.0)
add_textbox(slide29, "Team", pl29 + 0.2, pt29 + 0.45, pw29 - 0.4, 0.45,
            font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
# + Add button
add_rect(slide29, pl29 + pw29 - 1.2, pt29 + 0.42, 0.98, 0.38, fill=PURPLE, radius=True)
add_textbox(slide29, "+ Add", pl29 + pw29 - 1.2, pt29 + 0.44, 0.98, 0.3,
            font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Section labels
add_textbox(slide29, "DMs", pl29 + 0.22, pt29 + 1.0, pw29 - 0.44, 0.3,
            font_size=11, bold=True, color=PURPLE)

team_members = [
    ("J. Brown", "DM", PURPLE, True),
    ("S. Kim", "DM", PURPLE, True),
    ("Employees", None, None, False),
    ("M. Lee", "Employee", MED_GRAY, True),
    ("P. Davis", "Employee", MED_GRAY, True),
]
ty29 = pt29 + 1.3
for tm_name, tm_role, tm_col, tm_is_user in team_members:
    if tm_role is None:
        # section header
        add_textbox(slide29, tm_name, pl29 + 0.22, ty29 + 0.05, pw29 - 0.44, 0.3,
                    font_size=11, bold=True, color=PURPLE)
        ty29 += 0.38
        continue
    add_rect(slide29, pl29 + 0.15, ty29, pw29 - 0.3, 0.75, fill=DARK2, radius=True)
    add_textbox(slide29, tm_name, pl29 + 0.28, ty29 + 0.08, pw29 - 0.8, 0.32,
                font_size=11, bold=True, color=WHITE)
    # Role badge
    rb = add_rect(slide29, pl29 + 0.28, ty29 + 0.4, 1.0, 0.26, fill=tm_col, radius=True)
    add_textbox(slide29, tm_role, pl29 + 0.28, ty29 + 0.41, 1.0, 0.24,
                font_size=8, color=WHITE, align=PP_ALIGN.CENTER)
    # Edit button
    eb = add_rect(slide29, pl29 + pw29 - 1.1, ty29 + 0.2, 0.8, 0.35, fill=MED_GRAY, radius=True)
    add_textbox(slide29, "Edit", pl29 + pw29 - 1.1, ty29 + 0.22, 0.8, 0.28,
                font_size=9, color=WHITE, align=PP_ALIGN.CENTER)
    ty29 += 0.82

nav_bar(slide29, pl29, pt29 + ph29 - 0.95, pw29, 0.95, ["Home", "Tasks", "Flags", "Team", "More"])

add_slide_number(slide29, 29)

# ─── Slide 30 — Owner: Assigning Tasks ───────────────────────────────────────
slide30 = new_slide()
fill_bg(slide30, DARK_BG)
heading(slide30, "Owner: Assigning Tasks")
divider(slide30, 1.15)

bullet_list(slide30, [
    "Navigate to 'Tasks' in the nav bar.",
    "Tap '+ Assign Task' to create a new task.",
    "Assign to: Sales Director, Ops Manager, or DM.",
    "Set a due date for accountability.",
    "Track completion status directly in the Tasks view.",
], l=0.5, t=1.35, w=6.5, font_size=14, spacing=0.7)

# Phone mockup — task list
pl30, pt30, pw30, ph30 = phone_frame(slide30, l=7.55, t=0.55, w=4.95, h=6.0)
add_textbox(slide30, "Tasks", pl30 + 0.2, pt30 + 0.45, pw30 - 0.8, 0.45,
            font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(slide30, pl30 + pw30 - 1.25, pt30 + 0.42, 1.05, 0.38, fill=PURPLE, radius=True)
add_textbox(slide30, "+ Assign", pl30 + pw30 - 1.25, pt30 + 0.44, 1.05, 0.3,
            font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

tasks30 = [
    ("Submit Monthly Report", "J. Brown - DM", "Apr 16", RED, "Overdue"),
    ("New Staff Training", "S. Kim - DM", "Apr 20", ORANGE, "Pending"),
    ("Inventory Audit", "Ops Mgr", "Apr 14", GREEN, "Done"),
]
for ti30, (ttitle, tassignee, tdue, tcol, tstatus) in enumerate(tasks30):
    ty30 = pt30 + 1.05 + ti30 * 1.55
    add_rect(slide30, pl30 + 0.15, ty30, pw30 - 0.3, 1.3, fill=DARK2, radius=True)
    add_textbox(slide30, ttitle, pl30 + 0.28, ty30 + 0.08, pw30 - 0.55, 0.35,
                font_size=11, bold=True, color=WHITE)
    add_textbox(slide30, tassignee, pl30 + 0.28, ty30 + 0.42, pw30 - 0.55, 0.28,
                font_size=9, color=LIGHT_GRAY)
    # Status badge
    sb = add_rect(slide30, pl30 + 0.28, ty30 + 0.72, 1.2, 0.3, fill=tcol, radius=True)
    add_textbox(slide30, tstatus, pl30 + 0.28, ty30 + 0.73, 1.2, 0.26,
                font_size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide30, f"Due {tdue}", pl30 + pw30 - 1.4, ty30 + 0.72, 1.1, 0.28,
                font_size=8, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)
    if tstatus == "Done":
        add_textbox(slide30, "✓", pl30 + pw30 - 0.6, ty30 + 0.08, 0.4, 0.35,
                    font_size=16, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

nav_bar(slide30, pl30, pt30 + ph30 - 0.95, pw30, 0.95, ["Home", "Tasks", "Flags", "Team", "More"])

add_slide_number(slide30, 30)

# ─── Slide 31 — Owner: Store Inspection ──────────────────────────────────────
slide31 = new_slide()
fill_bg(slide31, DARK_BG)
heading(slide31, "Owner: Store Inspection (DM Visit)")
divider(slide31, 1.15)

bullet_list(slide31, [
    "Navigate to 'DM Visit' in the nav bar.",
    "Use this to log your store inspection visits.",
    "Record: visit date, store name, notes, and outcomes.",
    "Attach photos to document conditions.",
    "Completed inspections are stored for reference and accountability.",
], l=0.5, t=1.35, w=6.5, font_size=14, spacing=0.7)

# Phone mockup — DM Visit form
pl31, pt31, pw31, ph31 = phone_frame(slide31, l=7.55, t=0.55, w=4.95, h=6.0)
add_textbox(slide31, "DM Visit / Inspection", pl31 + 0.1, pt31 + 0.45, pw31 - 0.2, 0.45,
            font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

visit_fields = [
    ("Date of Visit", "Apr 15, 2026  📅"),
    ("Store Name", "Main St. Location  ▾"),
    ("DM Present", "J. Brown  ▾"),
]
for vfi, (vfl, vfv) in enumerate(visit_fields):
    vfy = pt31 + 1.1 + vfi * 0.9
    add_textbox(slide31, vfl, pl31 + 0.22, vfy, pw31 - 0.44, 0.26,
                font_size=9, color=LIGHT_GRAY)
    add_rect(slide31, pl31 + 0.22, vfy + 0.25, pw31 - 0.44, 0.48, fill=MED_GRAY, radius=True)
    add_textbox(slide31, vfv, pl31 + 0.36, vfy + 0.3, pw31 - 0.72, 0.35,
                font_size=10, color=WHITE)

# Notes area
notes_y = pt31 + 3.85
add_textbox(slide31, "Notes / Observations", pl31 + 0.22, notes_y, pw31 - 0.44, 0.26,
            font_size=9, color=LIGHT_GRAY)
add_rect(slide31, pl31 + 0.22, notes_y + 0.25, pw31 - 0.44, 1.0, fill=MED_GRAY, radius=True)
add_textbox(slide31, "Enter inspection notes...", pl31 + 0.36, notes_y + 0.38, pw31 - 0.72, 0.65,
            font_size=10, color=LIGHT_GRAY, italic=True)

# Submit button
add_rect(slide31, pl31 + 0.22, pt31 + 5.35, pw31 - 0.44, 0.5, fill=PURPLE, radius=True)
add_textbox(slide31, "Save Inspection", pl31 + 0.22, pt31 + 5.42, pw31 - 0.44, 0.38,
            font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_slide_number(slide31, 31)

# ─── Slide 32 — Owner: Map View ──────────────────────────────────────────────
slide32 = new_slide()
fill_bg(slide32, DARK_BG)
heading(slide32, "Owner: Map View")
divider(slide32, 1.15)

bullet_list(slide32, [
    "Navigate to 'Map' in the nav bar.",
    "View store locations and real-time GPS check-ins on a map.",
    "See which employees are currently clocked in and their locations.",
    "Useful for verifying on-site presence across multiple locations.",
], l=0.5, t=1.35, w=6.5, font_size=14, spacing=0.72)

# Phone mockup — Map view
pl32, pt32, pw32, ph32 = phone_frame(slide32, l=7.55, t=0.55, w=4.95, h=6.0)
add_textbox(slide32, "Map View", pl32 + 0.2, pt32 + 0.45, pw32 - 0.4, 0.45,
            font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Map background (simple grid)
add_rect(slide32, pl32 + 0.05, pt32 + 1.0, pw32 - 0.1, ph32 - 1.9,
         fill=RGBColor(0x1e, 0x3a, 0x5f))
# Grid lines (roads)
for gi in range(5):
    gy_line = pt32 + 1.0 + gi * 0.8
    add_rect(slide32, pl32 + 0.05, gy_line, pw32 - 0.1, 0.04,
             fill=RGBColor(0x2d, 0x4f, 0x7c))
for gi in range(4):
    gx_line = pl32 + 0.05 + gi * 1.2
    add_rect(slide32, gx_line, pt32 + 1.0, 0.04, ph32 - 1.9,
             fill=RGBColor(0x2d, 0x4f, 0x7c))

# Pin markers
pins = [(pl32 + 1.5, pt32 + 2.0, "Store A"), (pl32 + 3.2, pt32 + 3.0, "Store B")]
for px, py, plabel in pins:
    # Pin circle
    pc = slide32.shapes.add_shape(
        9, Inches(px), Inches(py), Inches(0.48), Inches(0.48)
    )
    pc.fill.solid()
    pc.fill.fore_color.rgb = PURPLE
    pc.line.color.rgb = WHITE
    pc.line.width = Pt(1.5)
    add_textbox(slide32, "📍", px - 0.05, py - 0.05, 0.55, 0.5,
                font_size=16, align=PP_ALIGN.CENTER)

    # Info card
    ic = add_rect(slide32, px - 0.1, py + 0.5, 1.6, 0.7, fill=DARK2, radius=True)
    add_textbox(slide32, plabel, px, py + 0.56, 1.4, 0.28,
                font_size=10, bold=True, color=WHITE)
    add_textbox(slide32, "2 clocked in", px, py + 0.82, 1.4, 0.28,
                font_size=9, color=GREEN)

nav_bar(slide32, pl32, pt32 + ph32 - 0.95, pw32, 0.95, ["Home", "Map", "Tasks", "Team", "More"])

add_slide_number(slide32, 32)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SECTION 7: RDM
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

# Slide 33 — RDM Section Header
section_header_slide("RDM", "Field Manager Pro")

# ─── Slide 34 — RDM: Overview ────────────────────────────────────────────────
slide34 = new_slide()
fill_bg(slide34, DARK_BG)
heading(slide34, "RDM: Overview & Responsibilities")
divider(slide34, 1.15)

add_textbox(slide34,
    "The RDM role provides a cross-organization view — spanning all organizations in the system.",
    0.4, 1.3, 12.0, 0.55, font_size=15, color=LIGHT_GRAY, italic=True)

bullet_list(slide34, [
    "No clock in/out function for RDM role.",
    "Can VIEW all staff schedules across all organizations.",
    "Can ASSIGN tasks to DMs across all organizations.",
    "No administrative control — view and task assignment only.",
], l=0.5, t=1.95, w=12.2, font_size=14, spacing=0.65)

add_textbox(slide34, "What the RDM can ACCESS:", 0.4, 4.5, 5.5, 0.38,
            font_size=15, bold=True, color=WHITE)

access_items = [
    ("📅", "Staff Schedule", "View Only"),
    ("✅", "Tasks", "View + Assign to DMs"),
    ("📊", "Dashboard", "View Only"),
]
for ai, (icon, aname, alevel) in enumerate(access_items):
    ay = 5.0 + ai * 0.72
    add_rect(slide34, 0.4, ay, 9.0, 0.62, fill=DARK2, radius=True)
    add_textbox(slide34, icon, 0.55, ay + 0.1, 0.5, 0.42, font_size=18)
    add_textbox(slide34, aname, 1.15, ay + 0.12, 3.5, 0.38,
                font_size=14, bold=True, color=WHITE)
    pill = add_rect(slide34, 4.8, ay + 0.14, 2.5, 0.34, fill=PURPLE, radius=True)
    add_textbox(slide34, alevel, 4.8, ay + 0.15, 2.5, 0.3,
                font_size=10, color=WHITE, align=PP_ALIGN.CENTER)

add_slide_number(slide34, 34)

# ─── Slide 35 — RDM: Assigning Tasks to DMs ──────────────────────────────────
slide35 = new_slide()
fill_bg(slide35, DARK_BG)
heading(slide35, "RDM: Assigning Tasks to DMs")
divider(slide35, 1.15)

steps35 = [
    ("1.", "Navigate to 'Tasks' in the nav bar."),
    ("2.", "Tap '+ Assign Task'."),
    ("3.", "Select any DM from the assignee list (cross-org visibility)."),
    ("4.", "Set a title, description, and due date."),
    ("5.", "Tap 'Assign Task' — the DM receives an email notification."),
    ("6.", "Track completion status in the Tasks view."),
]
for i, (num, txt) in enumerate(steps35):
    y = 1.35 + i * 0.65
    add_rect(slide35, 0.4, y, 0.38, 0.48, fill=PURPLE, radius=True)
    add_textbox(slide35, num, 0.42, y + 0.04, 0.35, 0.38,
                font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide35, txt, 0.9, y + 0.07, 5.9, 0.48, font_size=13, color=WHITE)

# Modal mockup showing cross-org DM list
pl35, pt35, pw35, ph35 = phone_frame(slide35, l=7.55, t=0.55, w=4.95, h=6.0)
add_rect(slide35, pl35, pt35 + 0.3, pw35, ph35 - 0.3, fill=RGBColor(0x00, 0x00, 0x00))
modal35_t = pt35 + 1.2
add_rect(slide35, pl35, modal35_t, pw35, ph35 - 1.2, fill=DARK2, radius=True)
add_rect(slide35, pl35 + pw35/2 - 0.3, modal35_t + 0.1, 0.6, 0.08, fill=MED_GRAY, radius=True)
add_textbox(slide35, "Assign Task", pl35 + 0.2, modal35_t + 0.25, pw35 - 0.4, 0.38,
            font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Assign To field with cross-org DMs
add_textbox(slide35, "Assign To", pl35 + 0.22, modal35_t + 0.7, pw35 - 0.44, 0.25,
            font_size=9, color=LIGHT_GRAY)
add_rect(slide35, pl35 + 0.22, modal35_t + 0.93, pw35 - 0.44, 0.48, fill=MED_GRAY, radius=True)

# Dropdown open showing DMs from diff orgs
dm_list = [
    ("J. Brown", "Org: West Region"),
    ("S. Kim", "Org: East Region"),
    ("A. Patel", "Org: South Region"),
]
drop_y = modal35_t + 1.45
add_rect(slide35, pl35 + 0.22, drop_y, pw35 - 0.44, len(dm_list) * 0.65 + 0.1, fill=MED_GRAY)
for dmi, (dmname, dmorg) in enumerate(dm_list):
    item_y = drop_y + 0.05 + dmi * 0.65
    item_bg = PURPLE if dmi == 0 else DARK2
    add_rect(slide35, pl35 + 0.22, item_y, pw35 - 0.44, 0.6, fill=item_bg)
    add_textbox(slide35, dmname, pl35 + 0.35, item_y + 0.04, pw35 - 0.7, 0.28,
                font_size=10, bold=True, color=WHITE)
    add_textbox(slide35, dmorg, pl35 + 0.35, item_y + 0.3, pw35 - 0.7, 0.25,
                font_size=8, color=LIGHT_GRAY)

# Assign button
ab_y = modal35_t + ph35 - 1.45
add_rect(slide35, pl35 + 0.22, ab_y, pw35 - 0.44, 0.52, fill=PURPLE, radius=True)
add_textbox(slide35, "Assign Task", pl35 + 0.22, ab_y + 0.06, pw35 - 0.44, 0.4,
            font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_slide_number(slide35, 35)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SECTION 8: CLOSING
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

# ─── Slide 36 — Support & Contact ────────────────────────────────────────────
slide36 = new_slide()
fill_bg(slide36, DARK_BG)
heading(slide36, "Support & Contact")
divider(slide36, 1.15)

add_textbox(slide36, "If you experience any issues with the app:", 0.4, 1.3, 12.0, 0.45,
            font_size=16, bold=True, color=WHITE)

support_steps = [
    ("Step 1", "Take a screenshot of the error message or unexpected behavior."),
    ("Step 2", "Email your Owner with a clear description and attach the screenshot."),
    ("Step 3", "The Owner will forward the issue to the developer for resolution."),
]
for si, (stitle, sdesc) in enumerate(support_steps):
    sy = 1.9 + si * 1.15
    add_rect(slide36, 0.4, sy, 1.6, 0.85, fill=PURPLE, radius=True)
    add_textbox(slide36, stitle, 0.4, sy + 0.2, 1.6, 0.45,
                font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide36, 2.1, sy, 10.0, 0.85, fill=DARK2, radius=True)
    add_textbox(slide36, sdesc, 2.3, sy + 0.2, 9.7, 0.45, font_size=13, color=WHITE)

# Warning
warn36 = add_rect(slide36, 0.4, 5.4, 12.5, 0.85, fill=RGBColor(0x7f, 0x1d, 0x1d), radius=True)
add_textbox(slide36, "🚫  Do NOT share your login credentials or attempt to clock in for another employee.",
            0.6, 5.55, 12.2, 0.55, font_size=13, bold=True, color=RED)

add_textbox(slide36, "fieldmanagerpro.app", 0.4, 6.55, 12.5, 0.45,
            font_size=16, color=PURPLE, align=PP_ALIGN.CENTER)

add_slide_number(slide36, 36)

# ─── Slide 37 — Thank You ────────────────────────────────────────────────────
slide37 = new_slide()
fill_bg(slide37, PURPLE)
add_rect(slide37, 0, 0, 13.33, 0.5, fill=DARK2)
add_rect(slide37, 0, 7.0, 13.33, 0.5, fill=DARK2)

add_textbox(slide37, "Welcome to", 0.5, 2.3, 12.3, 0.65,
            font_size=28, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide37, "Field Manager Pro", 0.5, 2.95, 12.3, 1.2,
            font_size=58, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide37, "You now have everything you need to get started.", 0.5, 4.3, 12.3, 0.65,
            font_size=22, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide37, "fieldmanagerpro.app", 0.5, 5.2, 12.3, 0.5,
            font_size=18, color=WHITE, align=PP_ALIGN.CENTER)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SAVE
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
out_path = "/Users/shaungephart/Desktop/Claude Projects/field-manager-pro/Field_Manager_Pro_User_Guide.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Total slides: {len(prs.slides)}")
